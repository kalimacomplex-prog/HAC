import asyncio
import base64
import glob as glob_module
import hashlib
import inspect
import io
import json
import os
import re
import secrets
import shutil
import subprocess
import sys
import time
import unicodedata
import zipfile
from datetime import datetime, timedelta
from typing import List

import httpx
from bson import ObjectId
from fastapi import APIRouter, Depends, HTTPException, Request
from pydantic import BaseModel

from ..auth import get_current_user
from ..config import settings
from ..database import pipelines_col, ai_agents_col, studio_automations_col, studio_runs_col, processes_col, jobs_col
from ..models.studio import (
    AutomationCreate, AutomationUpdate, AutomationOut, AutomationRunOut,
    StepResult, TriggerType, AutomationStep,
)
from .ai_agents import call_ai

router = APIRouter(prefix="/studio", tags=["studio"])


# ─── Helpers ──────────────────────────────────────────────────────

def _doc_to_out(doc: dict, base_url: str = "") -> AutomationOut:
    trigger = doc.get("trigger", {})
    webhook_url = ""
    if trigger.get("type") == "webhook" and trigger.get("webhook_token"):
        webhook_url = f"{base_url}/studio/webhook/{trigger['webhook_token']}"
    created = doc.get("created_at", "")
    return AutomationOut(
        id=doc["_id"],
        name=doc["name"],
        description=doc.get("description", ""),
        trigger=trigger,
        steps=doc.get("steps", []),
        active=doc.get("active", True),
        created_at=created.isoformat() if isinstance(created, datetime) else str(created),
        webhook_url=webhook_url,
        agent_id=doc.get("agent_id", ""),
    )


async def _upsert_linked_process(automation_id: str, user_id: str, name: str, description: str, agent_id: str, schedule: str):
    """Cria ou atualiza o Process vinculado a esta automação Studio."""
    now = datetime.utcnow()
    stub_script = f"# HAC Studio Automation\n# ID: {automation_id}\n# Gerenciado pelo HAC Studio Builder"
    existing = await processes_col.find_one({"studio_automation_id": automation_id, "user_id": user_id})
    if existing:
        await processes_col.update_one(
            {"_id": existing["_id"]},
            {"$set": {
                "name": name,
                "description": description,
                "agent_id": agent_id or None,
                "schedule": schedule or None,
                "updated_at": now,
            }}
        )
    else:
        doc = {
            "_id": str(ObjectId()),
            "user_id": user_id,
            "name": name,
            "description": description,
            "script": stub_script,
            "timeout_seconds": 300,
            "agent_id": agent_id or None,
            "schedule": schedule or None,
            "studio_automation_id": automation_id,
            "created_at": now,
            "updated_at": now,
        }
        await processes_col.insert_one(doc)


def _run_doc_to_out(doc: dict) -> AutomationRunOut:
    started = doc.get("started_at", "")
    finished = doc.get("finished_at")
    return AutomationRunOut(
        id=doc["_id"],
        automation_id=doc["automation_id"],
        automation_name=doc.get("automation_name", ""),
        status=doc["status"],
        trigger_type=doc.get("trigger_type", "manual"),
        input=doc.get("input", ""),
        steps_result=[StepResult(**s) for s in doc.get("steps_result", [])],
        output=doc.get("output", ""),
        started_at=started.isoformat() if isinstance(started, datetime) else str(started),
        finished_at=finished.isoformat() if isinstance(finished, datetime) else (str(finished) if finished else None),
        duration_ms=doc.get("duration_ms", 0),
    )


def _resolve_var_expr(expr: str, ctx: dict):
    """Resolve uma expressão tipo 'nome' ou 'nome[indice]' (indexação em lista/objeto
    JSON guardado numa variável — indice pode ser número, "chave entre aspas" ou o
    nome de outra variável, e dá pra encadear mais de um: nome[i][j]). Usada tanto
    por {var} / {var[indice]} em _sub() quanto pelas palavras soltas do passo 'log'.
    Retorna None se a variável base (antes do primeiro colchete) não existe — quem
    chama decide o que fazer nesse caso (não substituir, manter o texto original etc.),
    e uma string "expr (erro: ...)" se o índice/chave não bater, em vez de derrubar
    a automação inteira por causa disso."""
    m = re.match(r'^([A-Za-z_]\w*)((?:\[[^\]]*\])*)$', expr)
    if not m:
        return None
    base, brackets = m.group(1), m.group(2)
    variables = ctx.get("vars", {})
    if base == "output":
        value = ctx.get("output", "")
    elif base == "input":
        value = ctx.get("input", "")
    elif base in variables:
        value = variables[base]
    else:
        return None

    indices = re.findall(r'\[([^\]]*)\]', brackets)
    if not indices:
        return str(value)

    try:
        for idx_raw in indices:
            idx_raw = idx_raw.strip()
            if isinstance(value, str):
                value = json.loads(value)
            if idx_raw.startswith('"') and idx_raw.endswith('"'):
                key = idx_raw[1:-1]
            elif idx_raw in variables:
                key = variables[idx_raw]
            else:
                key = idx_raw
            if isinstance(value, list):
                value = value[int(key)]
            elif isinstance(value, dict):
                value = value.get(str(key))
            else:
                raise TypeError(f"'{value}' não é uma lista nem um objeto JSON")
        return str(value)
    except Exception as e:
        return f"{expr} (erro: {e})"


def _sub(text: str, ctx: dict) -> str:
    """Substitui {input}, {output}, {varname} e {varname[indice]} (ver _resolve_var_expr)
    no texto. Trecho não reconhecido (nome de variável que não existe) fica como estava —
    não mexe em chaves que não são referência a variável (ex: JSON literal no texto)."""
    text = str(text)

    def _replace(m):
        resolved = _resolve_var_expr(m.group(1), ctx)
        return m.group(0) if resolved is None else resolved

    return re.sub(r'\{([A-Za-z_]\w*(?:\[[^\]]*\])*)\}', _replace, text)


def _store(result: str, variable_name: str, ctx: dict):
    """Salva resultado no contexto. Se variable_name, salva em vars; sempre atualiza output."""
    ctx["output"] = result
    if variable_name:
        ctx["vars"][variable_name] = result


def _slugify_var(name: str) -> str:
    """Transforma um cabeçalho de coluna (ex: 'Preço Unitário') num nome de variável
    previsível (ex: 'preco_unitario') — usado por 'criar variáveis por coluna' em
    Ler Excel/Ler CSV. Sem acento, minúsculo, só letras/números/underscore."""
    s = unicodedata.normalize("NFKD", str(name)).encode("ascii", "ignore").decode("ascii")
    s = re.sub(r"[^a-zA-Z0-9]+", "_", s).strip("_").lower()
    return s or "coluna"


def _find_tesseract_binary():
    """Acha o executável do tesseract mesmo quando ele está instalado mas NÃO está
    no PATH do processo — comum no Windows logo após instalar via winget/instalador,
    já que o PATH do sistema só é atualizado pra processos novos, não pros que já
    estavam rodando. Sem isso, 'já está instalado' e 'está no PATH' viravam a mesma
    pergunta, e não são."""
    found = shutil.which("tesseract")
    if found:
        return found
    candidates = []
    if sys.platform == "win32":
        for base in (os.environ.get("PROGRAMFILES", ""), os.environ.get("PROGRAMFILES(X86)", ""),
                     os.environ.get("LOCALAPPDATA", "")):
            if base:
                candidates.append(os.path.join(base, "Tesseract-OCR", "tesseract.exe"))
    else:
        candidates += ["/usr/bin/tesseract", "/usr/local/bin/tesseract", "/opt/homebrew/bin/tesseract"]
    for c in candidates:
        if c and os.path.exists(c):
            return c
    return None


def _tesseract_installed() -> bool:
    return _find_tesseract_binary() is not None


def _find_poppler_bin_dir():
    """Mesma ideia de _find_tesseract_binary, mas devolve a PASTA (pdf2image pede o
    diretório, não o executável) — poppler no Windows costuma vir de um zip portátil
    ou do winget, sem entrar no PATH automaticamente."""
    found = shutil.which("pdftoppm")
    if found:
        return os.path.dirname(found)
    candidates = []
    if sys.platform == "win32":
        for base in (os.environ.get("PROGRAMFILES", ""), os.environ.get("LOCALAPPDATA", "")):
            if base and os.path.isdir(base):
                for entry in os.listdir(base):
                    if "poppler" in entry.lower():
                        for sub in ("Library/bin", "bin"):
                            cand = os.path.join(base, entry, *sub.split("/"))
                            if os.path.exists(os.path.join(cand, "pdftoppm.exe")):
                                candidates.append(cand)
    else:
        candidates += ["/usr/bin", "/usr/local/bin", "/opt/homebrew/bin"]
    for c in candidates:
        if c and os.path.exists(os.path.join(c, "pdftoppm" + (".exe" if sys.platform == "win32" else ""))):
            return c
    return None


def _poppler_installed() -> bool:
    return _find_poppler_bin_dir() is not None


def _ensure_native_binary(name: str, check, winget_id: str, apt_pkg: str, brew_pkg: str) -> None:
    """Garante que um binário nativo que NÃO é instalável via pip (tesseract, poppler)
    existe na máquina, tentando o instalador certo pro SO do agente antes de desistir —
    winget no Windows, apt/dnf/yum/pacman no Linux (com sudo -n, nunca trava esperando
    senha), brew no macOS. `check()` diz se já está instalado. Se a instalação
    automática não funcionar (falta admin/sudo, sem gerenciador suportado etc.),
    levanta um erro explicando como instalar manualmente em vez de travar sem dizer o
    que fazer."""
    if check():
        return
    print(f"'{name}' ausente no agente — tentando instalar automaticamente...")
    try:
        if sys.platform == "win32" and shutil.which("winget"):
            subprocess.run(
                ["winget", "install", "--id", winget_id, "-e", "--silent",
                 "--accept-package-agreements", "--accept-source-agreements"],
                capture_output=True, text=True, timeout=300, stdin=subprocess.DEVNULL,
            )
        elif sys.platform == "darwin" and shutil.which("brew"):
            subprocess.run(["brew", "install", brew_pkg], capture_output=True, text=True,
                            timeout=300, stdin=subprocess.DEVNULL)
        elif shutil.which("apt-get"):
            subprocess.run(["sudo", "-n", "apt-get", "update", "-q"], capture_output=True, text=True,
                            timeout=120, stdin=subprocess.DEVNULL)
            subprocess.run(["sudo", "-n", "apt-get", "install", "-y", apt_pkg], capture_output=True, text=True,
                            timeout=300, stdin=subprocess.DEVNULL)
        elif shutil.which("dnf"):
            subprocess.run(["sudo", "-n", "dnf", "install", "-y", apt_pkg], capture_output=True, text=True,
                            timeout=300, stdin=subprocess.DEVNULL)
        elif shutil.which("yum"):
            subprocess.run(["sudo", "-n", "yum", "install", "-y", apt_pkg], capture_output=True, text=True,
                            timeout=300, stdin=subprocess.DEVNULL)
        elif shutil.which("pacman"):
            subprocess.run(["sudo", "-n", "pacman", "-Sy", "--noconfirm", apt_pkg], capture_output=True, text=True,
                            timeout=300, stdin=subprocess.DEVNULL)
    except Exception:
        pass
    if not check():
        raise Exception(
            f"'{name}' não está instalado no agente e a instalação automática não funcionou "
            f"(pode faltar permissão de admin/sudo, ou nenhum gerenciador de pacotes suportado foi "
            f"encontrado nesta máquina). Instale manualmente e reinicie o agente: "
            f"Windows -> winget install --id {winget_id} -e ; "
            f"Linux -> sudo apt install {apt_pkg} (ou dnf/yum/pacman equivalente) ; "
            f"macOS -> brew install {brew_pkg}."
        )


# Nome do módulo Python (o que aparece em "No module named 'X'") -> nome do pacote pip,
# só para os casos em que os dois divergem. Usado pelo auto-instalador abaixo.
_MODULE_TO_PACKAGE = {
    "PIL": "Pillow",
    "cv2": "opencv-python-headless<5.0.0",
    "bs4": "beautifulsoup4",
    "yaml": "pyyaml",
    "dns": "dnspython",
    "jose": "python-jose[cryptography]",
    "docx": "python-docx",
    "pptx": "python-pptx",
    "whois": "python-whois",
    "slugify": "python-slugify",
    "validate_docbr": "validate-docbr",
    "Crypto": "pycryptodome",
}

# Módulos cuja falta NÃO deve disparar auto-instalação (dependências internas do
# projeto, stdlib mal detectada, etc.) — proteção contra tentar "pip install" algo
# que não existe no PyPI com esse nome.
_NO_AUTO_INSTALL = {"api", "worker", "bson", "fastapi", "pydantic", "starlette"}


def _pip_package_for_module(module_name: str) -> str:
    root = (module_name or "").split(".")[0]
    return _MODULE_TO_PACKAGE.get(root, root)


async def _try_pip_install(package: str) -> tuple:
    """Tenta instalar um pacote pip em runtime. Retorna (sucesso, mensagem)."""
    try:
        proc = await asyncio.create_subprocess_exec(
            sys.executable, "-m", "pip", "install", "--quiet", "--disable-pip-version-check", package,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)
        if proc.returncode == 0:
            return True, "instalado com sucesso"
        return False, (stderr or stdout).decode(errors="replace")[-800:]
    except asyncio.TimeoutError:
        return False, "timeout (mais de 120s) tentando instalar"
    except Exception as e:
        return False, str(e)


def _pix_tlv(id_: str, value: str) -> str:
    return f"{id_}{len(value):02d}{value}"


def _pix_crc16(payload: str) -> str:
    poly = 0x1021
    result = 0xFFFF
    for b in payload.encode():
        result ^= b << 8
        for _ in range(8):
            result = ((result << 1) ^ poly) & 0xFFFF if result & 0x8000 else (result << 1) & 0xFFFF
    return format(result, "04X")


def _build_pix_payload(key: str, name: str, city: str, amount: str = "") -> str:
    name = (name or "RECEBEDOR")[:25].upper()
    city = (city or "SAO PAULO")[:15].upper()
    merchant_account = _pix_tlv("26", _pix_tlv("00", "br.gov.bcb.pix") + _pix_tlv("01", key))
    payload = (
        _pix_tlv("00", "01")
        + merchant_account
        + _pix_tlv("52", "0000")
        + _pix_tlv("53", "986")
        + (_pix_tlv("54", f"{float(amount):.2f}") if amount else "")
        + _pix_tlv("58", "BR")
        + _pix_tlv("59", name)
        + _pix_tlv("60", city)
        + _pix_tlv("62", _pix_tlv("05", "***"))
    )
    payload += "6304"
    return payload + _pix_crc16(payload)


def _eval_condition(output: str, operator: str, value: str) -> bool:
    out_l, val_l = output.lower(), value.lower()
    if operator == "contains":       return val_l in out_l
    if operator == "not_contains":   return val_l not in out_l
    if operator == "equals":         return output.strip() == value.strip()
    if operator == "not_equals":     return output.strip() != value.strip()
    if operator == "starts_with":    return out_l.startswith(val_l)
    if operator == "ends_with":      return out_l.endswith(val_l)
    if operator == "is_empty":       return output.strip() == ""
    if operator == "not_empty":      return output.strip() != ""
    if operator == "greater_than":
        try: return float(output.strip()) > float(value.strip())
        except: return False
    if operator == "less_than":
        try: return float(output.strip()) < float(value.strip())
        except: return False
    return False


def _gen_browser_script(actions: list, engine: str, headless: bool, ctx: dict) -> str:
    """Gera script Python que será executado no agente."""
    subbed = [
        {
            "type": a.get("type", ""),
            "target": _sub(a.get("target", ""), ctx).replace("\\", "\\\\").replace('"', '\\"'),
            "value":  _sub(a.get("value",  ""), ctx).replace("\\", "\\\\").replace('"', '\\"'),
            "variable": a.get("variable", ""),
        }
        for a in actions
    ]

    if engine == "playwright":
        lines = [
            "import sys, os, json, subprocess",
            "sys.stdout.reconfigure(encoding='utf-8', errors='replace')",
            _HAC_ENSURE_PKG.strip("\n"),
            "_hac_ensure_pkg('playwright')",
            "from playwright.sync_api import sync_playwright",
            "_hac_ensure_playwright_chromium()",
            "_vars = {}",
            "with sync_playwright() as _pw:",
            f'    _br = _pw.chromium.launch(headless={str(headless)})',
            "    _pg = _br.new_page()",
            "    try:",
        ]
        for a in subbed:
            t, tgt, val, var = a["type"], a["target"], a["value"], a["variable"]
            indent = "        "
            if t == "open":
                url = tgt if tgt.startswith(("http://", "https://")) else f"https://{tgt}"
                lines += [f'{indent}_pg.goto("{url}", timeout=30000)', f'{indent}print("✓ Abriu: {url}")']
            elif t == "click":
                lines += [f'{indent}_pg.click("{tgt}", timeout=10000)', f'{indent}print("✓ Clicou: {tgt}")']
            elif t == "type":
                lines += [f'{indent}_pg.fill("{tgt}", "{val}", timeout=10000)', f'{indent}print("✓ Digitou em: {tgt}")']
            elif t == "extract":
                lines.append(f'{indent}_el = _pg.query_selector("{tgt}")')
                lines.append(f'{indent}_tx = (_el.get_attribute("{val}") if "{val}" else _el.inner_text()) if _el else ""')
                if var:
                    lines.append(f'{indent}_vars["{var}"] = str(_tx)')
                lines.append(f'{indent}print(f"✓ Extraiu: {{_tx}}")')
            elif t == "wait":
                ms = int(float(val or 1) * 1000)
                lines += [f'{indent}_pg.wait_for_timeout({ms})', f'{indent}print("✓ Aguardou {val or 1}s")']
            elif t == "screenshot":
                path = tgt or "/tmp/hac_screenshot.png"
                lines += [f'{indent}_pg.screenshot(path="{path}", full_page=True)', f'{indent}print("✓ Screenshot: {path}")']
            elif t == "close":
                lines += [f'{indent}_br.close()', f'{indent}print("✓ Browser fechado")']
        lines += [
            '        print("__VARS__:" + json.dumps(_vars))',
            "    finally:",
            "        try: _br.close()",
            "        except: pass",
        ]
    else:  # selenium
        lines = [
            "import sys, os, json, time, subprocess",
            "sys.stdout.reconfigure(encoding='utf-8', errors='replace')",
            _HAC_ENSURE_PKG.strip("\n"),
            "_hac_ensure_pkg('selenium')",
            "from selenium import webdriver",
            "from selenium.webdriver.common.by import By",
            "from selenium.webdriver.chrome.options import Options",
            "_vars = {}",
            "_opts = Options()",
        ]
        if headless:
            lines.append('_opts.add_argument("--headless=new")')
        lines += [
            '_opts.add_argument("--no-sandbox")',
            '_opts.add_argument("--disable-dev-shm-usage")',
            '_opts.add_argument("--disable-gpu")',
            "_dr = webdriver.Chrome(options=_opts)",
            "try:",
        ]
        for a in subbed:
            t, tgt, val, var = a["type"], a["target"], a["value"], a["variable"]
            indent = "    "
            if t == "open":
                url = tgt if tgt.startswith(("http://", "https://")) else f"https://{tgt}"
                lines += [f'{indent}_dr.get("{url}")', f'{indent}print("✓ Abriu: {url}")']
            elif t == "click":
                lines += [f'{indent}_dr.find_element(By.CSS_SELECTOR, "{tgt}").click()', f'{indent}print("✓ Clicou: {tgt}")']
            elif t == "type":
                lines += [f'{indent}_el = _dr.find_element(By.CSS_SELECTOR, "{tgt}")',
                          f'{indent}_el.clear(); _el.send_keys("{val}")',
                          f'{indent}print("✓ Digitou em: {tgt}")']
            elif t == "extract":
                lines.append(f'{indent}_el = _dr.find_element(By.CSS_SELECTOR, "{tgt}")')
                lines.append(f'{indent}_tx = _el.get_attribute("{val}") if "{val}" else _el.text')
                if var:
                    lines.append(f'{indent}_vars["{var}"] = str(_tx)')
                lines.append(f'{indent}print(f"✓ Extraiu: {{_tx}}")')
            elif t == "wait":
                lines += [f'{indent}time.sleep({float(val or 1)})', f'{indent}print("✓ Aguardou {val or 1}s")']
            elif t == "screenshot":
                path = tgt or "/tmp/hac_screenshot.png"
                lines += [f'{indent}_dr.save_screenshot("{path}")', f'{indent}print("✓ Screenshot: {path}")']
            elif t == "close":
                lines += [f'{indent}_dr.quit()', f'{indent}print("✓ Browser fechado")']
        lines += [
            '    print("__VARS__:" + json.dumps(_vars))',
            "finally:",
            "    try: _dr.quit()",
            "    except: pass",
        ]

    return "\n".join(lines)


async def _run_agent_script(script: str, agent_id: str, user_id: str,
                            timeout_seconds: int = 120, poll_seconds: int = 2,
                            name_prefix: str = "__studio_job_",
                            process_label: str = "Studio Step") -> str:
    """Despacha um script Python como job temporário (__studio_*) para o agente,
    aguarda a conclusão fazendo polling, sempre limpa o processo/job temporários
    (mesmo em falha/timeout) e retorna a saída bruta (stdout) do job."""
    now = datetime.utcnow()
    proc_id = str(ObjectId())
    job_id  = str(ObjectId())

    await processes_col.insert_one({
        "_id": proc_id, "user_id": user_id,
        "name": f"{name_prefix}{job_id[:8]}",
        "description": "Temporário — Studio Step",
        "script": script, "timeout_seconds": timeout_seconds,
        "agent_id": agent_id or None,
        "created_at": now, "updated_at": now,
    })
    await jobs_col.insert_one({
        "_id": job_id, "user_id": user_id,
        "process_id": proc_id, "process_name": process_label,
        "agent_id": agent_id or None,
        "status": "pending", "params": {}, "created_at": now,
    })

    # Aguarda conclusão (até timeout_seconds)
    job = None
    iterations = max(1, int(timeout_seconds / poll_seconds))
    for _ in range(iterations):
        await asyncio.sleep(poll_seconds)
        job = await jobs_col.find_one({"_id": job_id})
        if job and job["status"] in ("done", "failed", "cancelled"):
            break

    # Sempre deleta o processo/job temporários (mesmo em caso de falha/timeout)
    await processes_col.delete_one({"_id": proc_id})
    await jobs_col.delete_one({"_id": job_id})

    if not job or job["status"] != "done":
        err = (job or {}).get("error") or f"Timeout — agente não respondeu em {timeout_seconds}s"
        raise Exception(f"Execução no agente falhou: {err}")

    return job.get("output") or ""


async def _exec_browser_via_agent(actions: list, engine: str, headless: bool, ctx: dict,
                                   agent_id: str, user_id: str) -> tuple:
    """Gera script combinado e despacha como job para o agente, aguarda resultado."""
    script = _gen_browser_script(actions, engine, headless, ctx)
    raw_output = await _run_agent_script(
        script, agent_id, user_id, timeout_seconds=120,
        name_prefix="__studio_browser_", process_label="Studio Browser",
    )

    var_updates: dict = {}
    clean_lines = []
    for line in raw_output.splitlines():
        if line.startswith("__VARS__:"):
            try:
                var_updates = json.loads(line[9:])
            except Exception:
                pass
        else:
            clean_lines.append(line)

    return "\n".join(clean_lines), var_updates


# ── Sessões persistentes de navegador (CDP) ────────────────────────────────
#
# Em vez de despachar um único script combinado (modelo do step legado `browser`),
# cada ação vira um job pequeno e independente:
#   • "Abrir sessão" lança o Chromium real (binário do Playwright) como processo
#     DESTACADO, escutando numa porta de depuração remota (CDP) — esse processo
#     sobrevive ao fim do script que o lançou, ficando vivo na máquina do agente.
#   • As próximas ações (clicar, digitar, extrair, aguardar, screenshot) são scripts
#     curtos que se RECONECTAM nesse navegador via `connect_over_cdp`, fazem uma
#     única coisa e se desconectam — sem encerrar o processo remoto.
#   • "Fechar sessão" conecta, fecha as páginas, e então mata o processo pelo PID.
#
# Isso dá sessão persistente de verdade (sobrevive a steps não-navegador no meio
# do fluxo) sem precisar de NENHUMA mudança no worker — ele continua só recebendo
# e rodando scripts Python isolados, exatamente como já faz hoje.

_SESSION_LIFETIME_SECONDS = 30 * 60  # watchdog mata o navegador sozinho após esse tempo


# Embutido em TODOS os scripts de navegador (legado composto + sessão abrir/ação) —
# instala sozinho o pacote pip que falta (playwright/selenium) e, no caso do
# Playwright, também o próprio binário do Chromium quando ele ainda não foi baixado
# (equivalente a rodar "playwright install chromium" manualmente). Sem isso, faltar
# qualquer um dos dois derrubava o passo com ModuleNotFoundError/erro cru.
_HAC_ENSURE_PKG = """
import importlib as _hac_importlib

def _hac_ensure_pkg(pkg, import_name=None):
    import_name = import_name or pkg
    try:
        _hac_importlib.import_module(import_name)
        return
    except ModuleNotFoundError:
        pass
    print("Pacote '" + pkg + "' ausente no agente — instalando automaticamente...")
    _r = subprocess.run([sys.executable, "-m", "pip", "install", "--quiet", "--disable-pip-version-check", pkg],
                         capture_output=True, text=True, timeout=180)
    if _r.returncode != 0:
        raise ModuleNotFoundError(
            "Pacote '" + pkg + "' ausente no agente (interpretador " + sys.executable +
            ") e nao pode ser instalado automaticamente: " + (_r.stderr or _r.stdout)[-500:]
        )
    _hac_importlib.import_module(import_name)


def _hac_ensure_playwright_chromium():
    from playwright.sync_api import sync_playwright
    with sync_playwright() as _pw:
        _exe = _pw.chromium.executable_path
    if _exe and os.path.exists(_exe):
        return
    print("Chromium do Playwright ausente no agente — instalando automaticamente (pode demorar)...")
    _r = subprocess.run([sys.executable, "-m", "playwright", "install", "chromium"],
                         capture_output=True, text=True, timeout=300)
    if _r.returncode != 0:
        raise RuntimeError(
            "Nao foi possivel instalar o Chromium do Playwright automaticamente: " + (_r.stderr or _r.stdout)[-500:]
        )
"""


# Usado só por "Resolver captcha de imagem (OCR)". Tesseract "cru" praticamente não
# lê nada em captcha de imagem real (linhas riscando o texto, ruído de fundo,
# gradiente) — esse pré-processamento (escala de cinza, ampliação, desruído,
# binarização por Otsu, abertura morfológica pra apagar as linhas/pontos finos sem
# apagar o traço das letras) é o que faz a diferença entre ler algo e não ler nada.
_CAPTCHA_IMAGE_PREP = """
def _hac_read_captcha_text(img_path):
    _img = cv2.imread(img_path)
    _gray = cv2.cvtColor(_img, cv2.COLOR_BGR2GRAY)
    _gray = cv2.resize(_gray, None, fx=3, fy=3, interpolation=cv2.INTER_CUBIC)
    _gray = cv2.medianBlur(_gray, 3)
    _, _bin = cv2.threshold(_gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    _kernel = np.ones((2, 2), np.uint8)
    _clean = cv2.morphologyEx(_bin, cv2.MORPH_OPEN, _kernel)
    _cfg = "--psm 7 --oem 3 -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789"
    return pytesseract.image_to_string(_clean, config=_cfg).strip()
"""


_SELENIUM_CHROME_FINDER = """
def _find_chrome_browser():
    import shutil
    _names = ('chrome', 'google-chrome', 'google-chrome-stable', 'chromium', 'chromium-browser', 'chrome.exe')
    _cands = [shutil.which(_n) for _n in _names]
    if sys.platform == 'win32':
        for _base in (os.environ.get('PROGRAMFILES', ''), os.environ.get('PROGRAMFILES(X86)', ''), os.environ.get('LOCALAPPDATA', '')):
            if _base:
                _cands.append(os.path.join(_base, 'Google', 'Chrome', 'Application', 'chrome.exe'))
                _cands.append(os.path.join(_base, 'Chromium', 'Application', 'chrome.exe'))
        try:
            import winreg
            for _hive in (winreg.HKEY_LOCAL_MACHINE, winreg.HKEY_CURRENT_USER):
                try:
                    with winreg.OpenKey(_hive, r'SOFTWARE\\Microsoft\\Windows\\CurrentVersion\\App Paths\\chrome.exe') as _k:
                        _cands.append(winreg.QueryValue(_k, None))
                except OSError:
                    pass
        except Exception:
            pass
    elif sys.platform == 'darwin':
        _cands += ['/Applications/Google Chrome.app/Contents/MacOS/Google Chrome',
                   '/Applications/Chromium.app/Contents/MacOS/Chromium']
    else:
        _cands += ['/usr/bin/google-chrome', '/usr/bin/google-chrome-stable',
                   '/usr/bin/chromium', '/usr/bin/chromium-browser', '/snap/bin/chromium']
    for _c in _cands:
        if _c and os.path.exists(_c):
            return _c
    return None

"""


def _gen_session_open_script(session_name: str, target: str, headless: bool, ctx: dict,
                             engine: str = "playwright", profile_name: str = "") -> str:
    """Gera script que lança um Chrome/Chromium 'destacado' (sobrevive ao fim do
    processo que o lançou) escutando numa porta de depuração remota (CDP), e
    imprime `__SESSION__:{json com port/pid/user_data_dir}` para a API capturar.

    As duas engines ficam INDEPENDENTES: cada uma localiza e controla o navegador
    com sua própria biblioteca — Playwright nunca é importado para sessões Selenium
    e vice-versa.

    Se `profile_name` for informado, o profile do Chrome (user-data-dir) é uma pasta
    FIXA reaproveitada entre execuções (em vez de uma pasta temporária apagada ao
    fim da sessão) — assim extensões licenciadas instaladas manualmente uma vez
    (ex: para resolver captcha) continuam presentes nas próximas aberturas dessa
    mesma sessão nomeada. Extensões não funcionam em modo headless, então um
    profile persistente força headless=False."""
    url = _sub(target, ctx).strip()
    url_lit = json.dumps(url)
    name_lit = json.dumps(session_name)
    profile_name = (profile_name or "").strip()
    if profile_name:
        headless = False

    lines = ["import sys, os, json, socket, subprocess, tempfile, time, urllib.request",
             "sys.stdout.reconfigure(encoding='utf-8', errors='replace')",
             _HAC_ENSURE_PKG.strip("\n"), ""]

    if engine == "selenium":
        lines += [_SELENIUM_CHROME_FINDER.strip("\n")]
        lines += [
            "_exe = _find_chrome_browser()",
            "if not _exe:",
            "    print('__SESSION_ERROR__: Não foi possível localizar um navegador Chrome/Chromium instalado nesta máquina. '",
            "          'Instale o Google Chrome no agente e tente novamente.')",
            "    sys.exit(1)",
            "",
        ]
    else:
        lines += [
            "_hac_ensure_pkg('playwright')",
            "from playwright.sync_api import sync_playwright",
            "try:",
            "    _hac_ensure_playwright_chromium()",
            "except Exception as _e:",
            "    print('__SESSION_ERROR__: ' + str(_e))",
            "    sys.exit(1)",
            "with sync_playwright() as _pw:",
            "    _exe = _pw.chromium.executable_path",
            "",
        ]

    lines += [
        "_sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)",
        "_sock.bind(('127.0.0.1', 0))",
        "_port = _sock.getsockname()[1]",
        "_sock.close()",
        "",
    ]
    if profile_name:
        safe_profile = re.sub(r'[^A-Za-z0-9_-]+', '_', profile_name)[:80] or "default"
        lines += [
            "_udd = os.path.join(os.path.expanduser('~'), '.hac_browser_profiles', " + json.dumps(safe_profile) + ")",
            "os.makedirs(_udd, exist_ok=True)",
            "_udd_persistent = True",
        ]
    else:
        lines += [
            "_udd = tempfile.mkdtemp(prefix='hac_session_')",
            "_udd_persistent = False",
        ]
    lines += [
        "_args = [_exe, '--remote-debugging-port=%d' % _port, '--user-data-dir=' + _udd,",
        "         '--no-first-run', '--no-default-browser-check']",
        f"if {bool(headless)}:",
        "    _args.append('--headless=new')",
        "",
        "_kw = {'stdout': subprocess.DEVNULL, 'stderr': subprocess.DEVNULL, 'stdin': subprocess.DEVNULL}",
        "if sys.platform == 'win32':",
        "    _kw['creationflags'] = subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP",
        "else:",
        "    _kw['start_new_session'] = True",
        "_proc = subprocess.Popen(_args, **_kw)",
        "",
        "# watchdog destacado: garante que o navegador morre sozinho mesmo se a sessão nunca for fechada",
        "_wd = (",
        "    \"import time,sys,subprocess,shutil\\n\"",
        "    \"time.sleep(%d)\\n\"",
        "    \"if sys.platform == 'win32':\\n\"",
        "    \"    subprocess.run(['taskkill','/F','/T','/PID','%d'])\\n\"",
        "    \"else:\\n\"",
        "    \"    subprocess.run(['kill','-9','%d'])\\n\"",
        "    \"if not %s: shutil.rmtree(r'%s', ignore_errors=True)\\n\"",
        f") % ({_SESSION_LIFETIME_SECONDS}, _proc.pid, _proc.pid, _udd_persistent, _udd)",
        "subprocess.Popen([sys.executable, '-c', _wd], **_kw)",
        "",
        "# espera o endpoint de depuração remota (CDP) responder",
        "_ready = False",
        "for _ in range(60):",
        "    try:",
        "        urllib.request.urlopen('http://127.0.0.1:%d/json/version' % _port, timeout=1)",
        "        _ready = True",
        "        break",
        "    except Exception:",
        "        time.sleep(0.5)",
        "if not _ready:",
        "    print('__SESSION_ERROR__: navegador não respondeu na porta de depuração a tempo')",
        "    sys.exit(1)",
        "",
        f"_url = {url_lit}",
        "if _url:",
        "    _full = _url if _url.startswith(('http://', 'https://')) else 'https://' + _url",
        "    try:",
    ]
    if engine == "selenium":
        lines += [
            "        _hac_ensure_pkg('selenium')",
            "        from selenium import webdriver",
            "        _opts = webdriver.ChromeOptions()",
            "        _opts.add_experimental_option('debuggerAddress', '127.0.0.1:%d' % _port)",
            "        _dr = webdriver.Chrome(options=_opts)",
            "        try:",
            "            _dr.get(_full)",
            "        finally:",
            "            try: _dr.service.stop()",
            "            except Exception: pass",
        ]
    else:
        lines += [
            "        with sync_playwright() as _pw2:",
            "            _br = _pw2.chromium.connect_over_cdp('http://127.0.0.1:%d' % _port)",
            "            _bc = _br.contexts[0] if _br.contexts else _br.new_context()",
            "            _pg = _bc.pages[0] if _bc.pages else _bc.new_page()",
            "            _pg.goto(_full, timeout=30000)",
        ]
    lines += [
        "    except Exception as _e:",
        "        print('aviso: falha ao abrir URL inicial: ' + str(_e))",
        "",
        f"print('__SESSION__:' + json.dumps({{'port': _port, 'pid': _proc.pid, 'user_data_dir': _udd, "
        f"'session_name': {name_lit}, 'persistent': _udd_persistent}}))",
    ]
    return "\n".join(lines)


def _gen_session_action_script(action_type: str, port: int, target: str, value: str, ctx: dict,
                               engine: str = "playwright") -> str:
    """Gera script pequeno que se RECONECTA numa sessão já aberta — via CDP
    (Playwright) ou via `debuggerAddress` (Selenium) —, executa uma única ação
    e se desconecta SEM encerrar o navegador remoto (a sessão continua viva)."""
    tgt = _sub(target, ctx).replace("\\", "\\\\").replace('"', '\\"')
    val = _sub(value, ctx).replace("\\", "\\\\").replace('"', '\\"')
    indent = "    "

    # Ações de captcha (detectar/aguardar resolução manual/OCR em captcha de imagem)
    # usam esse trecho de JS pra checar se o token de resposta do reCAPTCHA/hCaptcha
    # já foi preenchido — é assim que dá pra saber que um humano resolveu o desafio
    # sem precisar de nenhum serviço pago de terceiro.
    _captcha_check_js = (
        "() => { const g = document.querySelector('textarea[name=\"g-recaptcha-response\"]'); "
        "const h = document.querySelector('textarea[name=\"h-captcha-response\"]'); "
        "return !!((g && g.value) || (h && h.value)); }"
    )

    if engine == "selenium":
        lines = [
            "import sys, os, time, re, subprocess, tempfile, shutil",
            "sys.stdout.reconfigure(encoding='utf-8', errors='replace')",
            _HAC_ENSURE_PKG.strip("\n"),
            "_hac_ensure_pkg('selenium')",
            "from selenium import webdriver",
            "from selenium.webdriver.common.by import By",
            "from selenium.webdriver.support.ui import WebDriverWait",
            "from selenium.webdriver.support import expected_conditions as EC",
            "",
            "# Detecta o tipo do seletor: prefixo explícito (xpath=/css=/id=/name=/class=/tag=/link=/partial_link=)",
            "# ou heurística (// , .. , ( -> XPath; caso contrário, CSS) — assim o mesmo campo aceita XPath,",
            "# CSS e os demais localizadores clássicos do Selenium.",
            "_BY_PREFIXES = {'xpath': By.XPATH, 'css': By.CSS_SELECTOR, 'id': By.ID, 'name': By.NAME,",
            "                'class': By.CLASS_NAME, 'tag': By.TAG_NAME, 'link': By.LINK_TEXT,",
            "                'partial_link': By.PARTIAL_LINK_TEXT}",
            "def _sel_locator(raw):",
            "    s = raw.strip()",
            "    m = re.match(r'^([a-zA-Z_]+)=(.*)$', s, re.S)",
            "    if m and m.group(1).lower() in _BY_PREFIXES:",
            "        return (_BY_PREFIXES[m.group(1).lower()], m.group(2).strip())",
            "    if s.startswith(('/', '..', '(')):",
            "        return (By.XPATH, s)",
            "    return (By.CSS_SELECTOR, s)",
            "",
        ]
        if action_type == "captcha_solve_image":
            lines += [
                inspect.getsource(_ensure_native_binary),
                inspect.getsource(_find_tesseract_binary),
                inspect.getsource(_tesseract_installed),
                "_hac_ensure_pkg('pytesseract')",
                "_hac_ensure_pkg('Pillow', 'PIL')",
                "_hac_ensure_pkg('opencv-python-headless<5.0.0', 'cv2')",
                "_hac_ensure_pkg('numpy')",
                "_ensure_native_binary('Tesseract OCR', _tesseract_installed, 'UB-Mannheim.TesseractOCR', 'tesseract-ocr', 'tesseract')",
                "import pytesseract, cv2, numpy as np",
                "from PIL import Image",
                "_tess_bin = _find_tesseract_binary()",
                "if _tess_bin: pytesseract.pytesseract.tesseract_cmd = _tess_bin",
                "",
                _CAPTCHA_IMAGE_PREP.strip("\n"),
            ]
        lines += [
            "_opts = webdriver.ChromeOptions()",
            f"_opts.add_experimental_option('debuggerAddress', '127.0.0.1:{port}')",
            "_dr = webdriver.Chrome(options=_opts)",
            "try:",
        ]
        if action_type == "click":
            # Tenta o clique normal; se outro elemento estiver por cima (ex: dropdown
            # de autocomplete cobrindo um botão de busca) ou o elemento não estiver
            # "interagível" pelas regras do WebDriver, cai para um clique via
            # JavaScript — que ignora essas checagens de sobreposição/visibilidade.
            lines += [
                f'{indent}_el = _dr.find_element(*_sel_locator("{tgt}"))',
                f'{indent}try:',
                f'{indent}    _el.click()',
                f'{indent}except Exception:',
                f'{indent}    _dr.execute_script("arguments[0].scrollIntoView({{block: \'center\'}}); arguments[0].click();", _el)',
                f'{indent}print("Clicou em: {tgt}")',
            ]
        elif action_type == "type":
            lines += [f'{indent}_el = _dr.find_element(*_sel_locator("{tgt}"))',
                      f'{indent}_el.clear(); _el.send_keys("{val}")',
                      f'{indent}print("Digitou em: {tgt}")']
        elif action_type == "extract":
            lines += [
                f'{indent}_el = _dr.find_element(*_sel_locator("{tgt}"))',
                f'{indent}_tx = _el.get_attribute("{val}") if "{val}" else _el.text',
                f'{indent}print(_tx)',
            ]
        elif action_type == "wait":
            if tgt:
                lines += [f'{indent}WebDriverWait(_dr, 30).until(EC.presence_of_element_located(_sel_locator("{tgt}")))',
                          f'{indent}print("Elemento apareceu: {tgt}")']
            else:
                lines += [f'{indent}time.sleep({float(val or 1)})', f'{indent}print("Aguardou {val or 1}s")']
        elif action_type == "screenshot":
            path = tgt or "hac_screenshot.png"
            lines += [f'{indent}_dr.save_screenshot("{path}")', f'{indent}print("Screenshot salvo em: {path}")']
        elif action_type == "captcha_detect":
            # Não existe forma gratuita de RESOLVER reCAPTCHA/hCaptcha automaticamente
            # (são desafios anti-bot de propósito) — mas dá pra DETECTAR qual apareceu,
            # pra decidir o que fazer a seguir (ex: pausar pra alguém resolver).
            lines += [
                f'{indent}_type = "nenhum"',
                f'{indent}for _fr in _dr.find_elements(By.TAG_NAME, "iframe"):',
                f'{indent}    _src = (_fr.get_attribute("src") or "").lower()',
                f'{indent}    if "recaptcha" in _src: _type = "recaptcha"; break',
                f'{indent}    if "hcaptcha" in _src: _type = "hcaptcha"; break',
                f'{indent}if _type == "nenhum" and "{tgt}":',
                f'{indent}    try:',
                f'{indent}        _dr.find_element(*_sel_locator("{tgt}"))',
                f'{indent}        _type = "imagem"',
                f'{indent}    except Exception:',
                f'{indent}        pass',
                f'{indent}print(_type)',
            ]
        elif action_type == "captcha_wait":
            # Espera alguém resolver o captcha NA TELA (a sessão precisa estar com
            # 'Onde executar: Agente' e sem headless pra um humano conseguir ver e
            # interagir) — sem seletor customizado, checa o token padrão do reCAPTCHA/
            # hCaptcha; com seletor, espera esse elemento aparecer (ex: mensagem de sucesso).
            check_js_lit = json.dumps(f"return ({_captcha_check_js})();")
            lines += [
                f'{indent}_timeout_s = float("{val}" or 120)',
                f'{indent}_deadline = time.time() + _timeout_s',
                f'{indent}_solved = False',
                f'{indent}while time.time() < _deadline:',
                f'{indent}    if "{tgt}":',
                f'{indent}        try:',
                f'{indent}            _dr.find_element(*_sel_locator("{tgt}"))',
                f'{indent}            _solved = True',
                f'{indent}            break',
                f'{indent}        except Exception:',
                f'{indent}            pass',
                f'{indent}    elif _dr.execute_script({check_js_lit}):',
                f'{indent}        _solved = True',
                f'{indent}        break',
                f'{indent}    time.sleep(1.5)',
                f'{indent}if not _solved:',
                f'{indent}    raise Exception("Tempo esgotado esperando a resolução manual do captcha (" + str(_timeout_s) + "s)")',
                f'{indent}print("Captcha resolvido")',
            ]
        elif action_type == "captcha_solve_image":
            # Só funciona bem em captcha de imagem "clássico" (texto distorcido) —
            # reCAPTCHA/hCaptcha não dá pra ler por OCR, foram feitos pra impedir isso.
            lines += [
                f'{indent}_el = _dr.find_element(*_sel_locator("{tgt}"))',
                f'{indent}_img_path = tempfile.mktemp(suffix=".png")',
                f'{indent}_el.screenshot(_img_path)',
                f'{indent}_captcha_text = _hac_read_captcha_text(_img_path)',
                f'{indent}print(_captcha_text)',
            ]
            if val:
                lines += [
                    f'{indent}if _captcha_text:',
                    f'{indent}    _el2 = _dr.find_element(*_sel_locator("{val}"))',
                    f'{indent}    _el2.clear(); _el2.send_keys(_captcha_text)',
                ]
        lines += [
            "finally:",
            # Importante: NUNCA chamar _dr.quit() aqui — numa sessão anexada via
            # debuggerAddress isso fecharia o navegador remoto. service.stop()
            # apenas encerra o processo local do chromedriver (desconecta).
            f"{indent}try: _dr.service.stop()",
            f"{indent}except Exception: pass",
        ]
        return "\n".join(lines)

    lines = [
        "import sys, os, re, time, tempfile, subprocess, shutil",
        "sys.stdout.reconfigure(encoding='utf-8', errors='replace')",
        _HAC_ENSURE_PKG.strip("\n"),
        "_hac_ensure_pkg('playwright')",
        "from playwright.sync_api import sync_playwright",
        "_hac_ensure_playwright_chromium()",
        "",
        "# Playwright já detecta XPath nativamente quando o seletor começa com '//' ou '..',",
        "# e suporta prefixos explícitos de engine (css=, xpath=, text=, role= etc). A única",
        "# lacuna é o XPath absoluto de barra única (ex: /html/body/...), comum ao copiar",
        "# 'Copy XPath' do DevTools — então normalizamos esse caso para 'xpath=...'.",
        "def _sel(raw):",
        "    s = raw.strip()",
        "    if re.match(r'^[a-zA-Z_]+=', s):",
        "        return s",
        "    if s.startswith('/') and not s.startswith('//'):",
        "        return 'xpath=' + s",
        "    return s",
        "",
    ]
    if action_type == "captcha_solve_image":
        lines += [
            inspect.getsource(_ensure_native_binary),
            inspect.getsource(_find_tesseract_binary),
            inspect.getsource(_tesseract_installed),
            "_hac_ensure_pkg('pytesseract')",
            "_hac_ensure_pkg('Pillow', 'PIL')",
            "_hac_ensure_pkg('opencv-python-headless<5.0.0', 'cv2')",
            "_hac_ensure_pkg('numpy')",
            "_ensure_native_binary('Tesseract OCR', _tesseract_installed, 'UB-Mannheim.TesseractOCR', 'tesseract-ocr', 'tesseract')",
            "import pytesseract, cv2, numpy as np",
            "from PIL import Image",
            "_tess_bin = _find_tesseract_binary()",
            "if _tess_bin: pytesseract.pytesseract.tesseract_cmd = _tess_bin",
            "",
            _CAPTCHA_IMAGE_PREP.strip("\n"),
        ]
    lines += [
        "with sync_playwright() as _pw:",
        f"    _br = _pw.chromium.connect_over_cdp('http://127.0.0.1:{port}')",
        "    _bc = _br.contexts[0] if _br.contexts else _br.new_context()",
        "    _pg = _bc.pages[0] if _bc.pages else _bc.new_page()",
    ]
    if action_type == "click":
        # Tenta o clique normal (com as checagens de actionability do Playwright);
        # se algo cobrir o elemento de forma persistente (ex: dropdown de
        # autocomplete sobre um botão de busca), tenta de novo com force=True,
        # que ignora a checagem de "recebe eventos de ponteiro".
        lines += [
            f'{indent}try:',
            f'{indent}    _pg.click(_sel("{tgt}"), timeout=15000)',
            f'{indent}except Exception:',
            f'{indent}    _pg.click(_sel("{tgt}"), timeout=5000, force=True)',
            f'{indent}print("Clicou em: {tgt}")',
        ]
    elif action_type == "type":
        lines += [f'{indent}_pg.fill(_sel("{tgt}"), "{val}", timeout=15000)', f'{indent}print("Digitou em: {tgt}")']
    elif action_type == "extract":
        lines += [
            f'{indent}_el = _pg.query_selector(_sel("{tgt}"))',
            f'{indent}_tx = (_el.get_attribute("{val}") if "{val}" else _el.inner_text()) if _el else ""',
            f'{indent}print(_tx)',
        ]
    elif action_type == "wait":
        if tgt:
            lines += [f'{indent}_pg.wait_for_selector(_sel("{tgt}"), timeout=30000)', f'{indent}print("Elemento apareceu: {tgt}")']
        else:
            ms = int(float(val or 1) * 1000)
            lines += [f'{indent}_pg.wait_for_timeout({ms})', f'{indent}print("Aguardou {val or 1}s")']
    elif action_type == "screenshot":
        path = tgt or "hac_screenshot.png"
        lines += [f'{indent}_pg.screenshot(path="{path}", full_page=True)', f'{indent}print("Screenshot salvo em: {path}")']
    elif action_type == "captcha_detect":
        # Não existe forma gratuita de RESOLVER reCAPTCHA/hCaptcha automaticamente
        # (são desafios anti-bot de propósito) — mas dá pra DETECTAR qual apareceu,
        # pra decidir o que fazer a seguir (ex: pausar pra alguém resolver).
        lines += [
            f'{indent}_type = "nenhum"',
            f'{indent}for _fr in _pg.query_selector_all("iframe"):',
            f'{indent}    _src = (_fr.get_attribute("src") or "").lower()',
            f'{indent}    if "recaptcha" in _src: _type = "recaptcha"; break',
            f'{indent}    if "hcaptcha" in _src: _type = "hcaptcha"; break',
            f'{indent}if _type == "nenhum" and "{tgt}" and _pg.query_selector(_sel("{tgt}")):',
            f'{indent}    _type = "imagem"',
            f'{indent}print(_type)',
        ]
    elif action_type == "captcha_wait":
        # Espera alguém resolver o captcha NA TELA (a sessão precisa estar com
        # 'Onde executar: Agente' e sem headless pra um humano conseguir ver e
        # interagir) — sem seletor customizado, checa o token padrão do reCAPTCHA/
        # hCaptcha; com seletor, espera esse elemento aparecer (ex: mensagem de sucesso).
        lines += [
            f'{indent}_timeout_s = float("{val}" or 120)',
            f'{indent}_deadline = time.time() + _timeout_s',
            f'{indent}_solved = False',
            f'{indent}while time.time() < _deadline:',
            f'{indent}    if "{tgt}" and _pg.query_selector(_sel("{tgt}")):',
            f'{indent}        _solved = True',
            f'{indent}        break',
            f'{indent}    elif not "{tgt}" and _pg.evaluate({json.dumps(_captcha_check_js)}):',
            f'{indent}        _solved = True',
            f'{indent}        break',
            f'{indent}    time.sleep(1.5)',
            f'{indent}if not _solved:',
            f'{indent}    raise Exception("Tempo esgotado esperando a resolução manual do captcha (" + str(_timeout_s) + "s)")',
            f'{indent}print("Captcha resolvido")',
        ]
    elif action_type == "captcha_solve_image":
        # Só funciona bem em captcha de imagem "clássico" (texto distorcido) —
        # reCAPTCHA/hCaptcha não dá pra ler por OCR, foram feitos pra impedir isso.
        lines += [
            f'{indent}_el = _pg.query_selector(_sel("{tgt}"))',
            f'{indent}if not _el: raise Exception("Elemento da imagem do captcha não encontrado: {tgt}")',
            f'{indent}_img_path = tempfile.mktemp(suffix=".png")',
            f'{indent}_el.screenshot(path=_img_path)',
            f'{indent}_captcha_text = _hac_read_captcha_text(_img_path)',
            f'{indent}print(_captcha_text)',
        ]
        if val:
            lines += [
                f'{indent}if _captcha_text:',
                f'{indent}    _pg.fill(_sel("{val}"), _captcha_text, timeout=10000)',
            ]
    lines.append(f"{indent}_br.close()")
    return "\n".join(lines)


def _gen_session_close_script(pid: int, port: int, user_data_dir: str, engine: str = "playwright",
                               persistent: bool = False) -> str:
    """Gera script que tenta fechar as páginas da sessão (melhor esforço, usando a
    MESMA biblioteca da engine escolhida — sem depender da outra) e então mata o
    processo do navegador pelo PID e remove o diretório temporário de perfil.

    Se `persistent` for True, o diretório NÃO é removido — é um profile fixo
    reaproveitado nas próximas aberturas dessa sessão (ver `_gen_session_open_script`)."""
    udd = user_data_dir.replace("\\", "\\\\").replace('"', '\\"')
    lines = ["import sys, subprocess, shutil", "sys.stdout.reconfigure(encoding='utf-8', errors='replace')"]

    if engine == "selenium":
        # Nada de melhor-esforço aqui: anexar via Selenium só para desconectar não
        # fecha nenhuma página, então vamos direto para o kill pelo PID abaixo —
        # mantém a engine Selenium livre de qualquer dependência do Playwright.
        pass
    else:
        lines += [
            "try:",
            "    from playwright.sync_api import sync_playwright",
            "    with sync_playwright() as _pw:",
            f"        _br = _pw.chromium.connect_over_cdp('http://127.0.0.1:{port}')",
            "        for _ctx in list(_br.contexts):",
            "            for _pg in list(_ctx.pages):",
            "                try: _pg.close()",
            "                except Exception: pass",
            "        _br.close()",
            "except Exception:",
            "    pass",
        ]

    lines += [
        "if sys.platform == 'win32':",
        f"    subprocess.run(['taskkill', '/F', '/T', '/PID', '{pid}'])",
        "else:",
        f"    subprocess.run(['kill', '-9', '{pid}'])",
    ]
    if not persistent:
        lines.append(f'shutil.rmtree("{udd}", ignore_errors=True)')
    lines.append('print("Sessão encerrada")')
    return "\n".join(lines)


async def _exec_step(step: dict, ctx: dict) -> str:
    """Executa um step e retorna o output. Atualiza ctx em tempo real."""
    t = step["type"]
    cfg = step.get("config", {})
    var = cfg.get("variable_name", "")

    # ── Controle de Fluxo ──────────────────────────────────

    if t == "wait":
        secs = float(cfg.get("seconds", 1))
        remaining = min(secs, 60)
        while remaining > 0:
            if await _is_cancelled(ctx):
                ctx["_cancelled"] = True
                raise Exception("Cancelada pelo usuário")
            chunk = min(0.5, remaining)
            await asyncio.sleep(chunk)
            remaining -= chunk
        return f"Aguardou {secs}s"

    if t == "comment":
        return f"# {cfg.get('text', '')}"

    if t == "log":
        # Tipo um print(a, "texto", b): cada palavra solta (sem aspas) é tratada como
        # nome de variável (ou output/input, opcionalmente indexada com [algo] — ver
        # _resolve_var_expr); texto literal precisa vir entre "aspas". Junta tudo com
        # espaço. Não passa por ctx['output'] — é só uma anotação de debug, não deve
        # mudar o {output} que o próximo passo recebe.
        raw = cfg.get("text", "")
        parts = []
        for m in re.finditer(r'"([^"]*)"|([A-Za-z_]\w*(?:\[[^\]]*\])*)|(\S+)', raw):
            if m.group(1) is not None:
                parts.append(m.group(1))
            elif m.group(2) is not None:
                resolved = _resolve_var_expr(m.group(2), ctx)
                parts.append(m.group(2) if resolved is None else resolved)
            else:
                parts.append(m.group(3))
        return " ".join(parts)

    # condition é tratado no loop principal (precisa controlar índice)
    # loop_count idem — retornamos placeholder aqui
    if t in ("condition", "loop_count"):
        return ""

    # ── Variáveis ──────────────────────────────────────────

    if t == "set_variable":
        result = _sub(cfg.get("value", ""), ctx)
        _store(result, var, ctx)
        return result

    if t == "calculate":
        expr = _sub(cfg.get("expression", "0"), ctx)
        try:
            import math
            safe = {"__builtins__": {}, "math": math, "abs": abs, "round": round,
                    "len": len, "str": str, "int": int, "float": float, "min": min, "max": max}
            safe.update(ctx.get("vars", {}))
            result = str(eval(expr, safe))
        except Exception as e:
            raise Exception(f"Erro ao calcular '{expr}': {e}")
        _store(result, var, ctx)
        return result

    # ── Arquivos ───────────────────────────────────────────

    if t == "read_file":
        path = _sub(cfg.get("file_path", ""), ctx)
        with open(path, "r", encoding="utf-8") as f:
            result = f.read()
        _store(result, var, ctx)
        return result

    if t == "write_file":
        path = _sub(cfg.get("file_path", ""), ctx)
        content = _sub(cfg.get("content", "{output}"), ctx)
        mode = "a" if cfg.get("append") else "w"
        os.makedirs(os.path.dirname(path), exist_ok=True) if os.path.dirname(path) else None
        with open(path, mode, encoding="utf-8") as f:
            f.write(content)
        result = f"Arquivo {'adicionado' if cfg.get('append') else 'escrito'}: {path}"
        _store(result, var, ctx)
        return result

    if t == "list_files":
        directory = _sub(cfg.get("directory", "."), ctx)
        pattern = cfg.get("pattern", "*")
        files = glob_module.glob(os.path.join(directory, pattern))
        result = json.dumps(files, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "delete_file":
        path = _sub(cfg.get("file_path", ""), ctx)
        os.remove(path)
        result = f"Arquivo deletado: {path}"
        _store(result, var, ctx)
        return result

    if t == "copy_file":
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        shutil.copy2(src, dst)
        result = f"Copiado: {src} → {dst}"
        _store(result, var, ctx)
        return result

    if t == "move_file":
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        shutil.move(src, dst)
        result = f"Movido: {src} → {dst}"
        _store(result, var, ctx)
        return result

    if t == "file_hash":
        path = _sub(cfg.get("file_path", ""), ctx)
        algo = cfg.get("hash_algo", "sha256")
        h = hashlib.new(algo)
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        result = h.hexdigest()
        _store(result, var, ctx)
        return result

    if t == "file_info":
        path = _sub(cfg.get("file_path", ""), ctx)
        st = os.stat(path)
        info = {
            "path": path, "size_bytes": st.st_size,
            "modified": datetime.fromtimestamp(st.st_mtime).isoformat(),
            "created": datetime.fromtimestamp(st.st_ctime).isoformat(),
            "is_dir": os.path.isdir(path),
        }
        result = json.dumps(info, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "search_in_files":
        directory = _sub(cfg.get("directory", "."), ctx)
        pattern = cfg.get("pattern", "*") or "*"
        needle = _sub(cfg.get("search", ""), ctx)
        matches = []
        for path in glob_module.glob(os.path.join(directory, "**", pattern), recursive=True):
            if not os.path.isfile(path):
                continue
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    for line_no, line in enumerate(f, 1):
                        if needle in line:
                            matches.append(f"{path}:{line_no}: {line.strip()[:200]}")
            except Exception:
                continue
        result = "\n".join(matches) if matches else "Nenhuma ocorrência encontrada."
        _store(result, var, ctx)
        return result

    if t == "convert_encoding":
        src = _sub(cfg.get("source_path", "") or cfg.get("file_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx) or src
        enc_from = cfg.get("encoding_from", "utf-8")
        enc_to = cfg.get("encoding_to", "utf-8")
        with open(src, "r", encoding=enc_from) as f:
            content = f.read()
        with open(dst, "w", encoding=enc_to) as f:
            f.write(content)
        result = f"Convertido {src} ({enc_from} → {enc_to}) em {dst}"
        _store(result, var, ctx)
        return result

    if t == "delete_folder":
        path = _sub(cfg.get("directory", "") or cfg.get("file_path", ""), ctx)
        shutil.rmtree(path)
        result = f"Pasta deletada: {path}"
        _store(result, var, ctx)
        return result

    if t == "ensure_dir":
        path = _sub(cfg.get("directory", "") or cfg.get("file_path", ""), ctx)
        os.makedirs(path, exist_ok=True)
        result = f"Diretório garantido: {path}"
        _store(result, var, ctx)
        return result

    # ── Compactação ─────────────────────────────────────────

    if t == "zip_files":
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        if not dst.lower().endswith(".zip"):
            dst += ".zip"
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zf:
            if os.path.isdir(src):
                for root, _dirs, files in os.walk(src):
                    for fn in files:
                        full = os.path.join(root, fn)
                        zf.write(full, os.path.relpath(full, src))
            else:
                zf.write(src, os.path.basename(src))
        result = f"Compactado em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "unzip_file":
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", "") or ".", ctx)
        os.makedirs(dst, exist_ok=True)
        with zipfile.ZipFile(src, "r") as zf:
            zf.extractall(dst)
        result = f"Descompactado em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "backup_folder":
        src = _sub(cfg.get("source_path", ""), ctx)
        dst_base = _sub(cfg.get("dest_path", ""), ctx)
        stamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        dst = f"{dst_base}_{stamp}"
        shutil.copytree(src, dst)
        result = f"Backup criado em: {dst}"
        _store(result, var, ctx)
        return result

    # ── Data & Hora ─────────────────────────────────────────

    if t == "date_diff":
        fmt_in = cfg.get("date_format_in", "%Y-%m-%d")
        d1 = datetime.strptime(_sub(cfg.get("date_value", ""), ctx), fmt_in)
        d2 = datetime.strptime(_sub(cfg.get("date_value2", ""), ctx), fmt_in)
        delta_seconds = (d2 - d1).total_seconds()
        unit = cfg.get("date_unit", "days")
        unit_map = {
            "seconds": delta_seconds, "minutes": delta_seconds / 60,
            "hours": delta_seconds / 3600, "days": delta_seconds / 86400,
            "weeks": delta_seconds / 604800,
        }
        result = str(round(unit_map.get(unit, delta_seconds / 86400), 2))
        _store(result, var, ctx)
        return result

    if t == "date_add":
        base = datetime.strptime(_sub(cfg.get("date_value", ""), ctx), cfg.get("date_format_in", "%Y-%m-%d"))
        unit = cfg.get("date_unit", "days")
        amount = int(cfg.get("date_amount", 0))
        kwargs = {unit: amount} if unit in ("seconds", "minutes", "hours", "days", "weeks") else {"days": amount}
        new_date = base + timedelta(**kwargs)
        result = new_date.strftime(cfg.get("date_format_out", "%Y-%m-%d"))
        _store(result, var, ctx)
        return result

    if t == "timezone_convert":
        from zoneinfo import ZoneInfo
        fmt_in = cfg.get("date_format_in", "%Y-%m-%d %H:%M:%S")
        naive = datetime.strptime(_sub(cfg.get("date_value", ""), ctx), fmt_in)
        aware = naive.replace(tzinfo=ZoneInfo(cfg.get("timezone_from", "UTC")))
        converted = aware.astimezone(ZoneInfo(cfg.get("timezone_to", "America/Sao_Paulo")))
        result = converted.strftime(cfg.get("date_format_out", "%Y-%m-%d %H:%M:%S"))
        _store(result, var, ctx)
        return result

    if t == "is_business_day":
        d = datetime.strptime(_sub(cfg.get("date_value", ""), ctx), cfg.get("date_format_in", "%Y-%m-%d"))
        result = "true" if d.weekday() < 5 else "false"
        _store(result, var, ctx)
        return result

    if t == "format_date":
        d = datetime.strptime(_sub(cfg.get("date_value", ""), ctx), cfg.get("date_format_in", "%Y-%m-%d"))
        result = d.strftime(cfg.get("date_format_out", "%d/%m/%Y"))
        _store(result, var, ctx)
        return result

    if t == "random_wait":
        import random
        lo = float(cfg.get("seconds", 1))
        hi = float(cfg.get("seconds_max", 3))
        if hi < lo:
            hi = lo
        secs = random.uniform(lo, min(hi, 60))
        remaining = secs
        while remaining > 0:
            if await _is_cancelled(ctx):
                ctx["_cancelled"] = True
                raise Exception("Cancelada pelo usuário")
            chunk = min(0.5, remaining)
            await asyncio.sleep(chunk)
            remaining -= chunk
        result = f"Aguardou {secs:.2f}s (aleatório entre {lo}s e {hi}s)"
        _store(result, var, ctx)
        return result

    if t == "call_automation":
        target_id = cfg.get("automation_id", "")
        if not target_id:
            raise Exception("Automação não selecionada")
        depth = ctx.get("call_depth", 0)
        if depth >= 5:
            raise Exception("Profundidade máxima de sub-fluxos (5) atingida — possível recursão")
        target = await studio_automations_col.find_one({"_id": target_id})
        if not target:
            raise Exception(f"Automação '{target_id}' não encontrada")
        sub_input = _sub(cfg.get("input_template", "{output}") or "{output}", ctx)
        sub_ctx = {
            "input": sub_input, "output": sub_input, "vars": {}, "sessions": ctx.get("sessions", {}),
            "agent_id": target.get("agent_id", ctx.get("agent_id", "")), "user_id": ctx.get("user_id", ""),
            "call_depth": depth + 1,
        }
        sub_results, sub_failed = await _exec_step_list(target.get("steps", []), sub_ctx)
        if sub_failed:
            last_err = next((r["error"] for r in reversed(sub_results) if r.get("error")), "erro desconhecido")
            raise Exception(f"Sub-fluxo '{target.get('name')}' falhou: {last_err}")
        result = sub_ctx["output"]
        _store(result, var, ctx)
        return result

    # ── Planilhas & Excel ───────────────────────────────────

    if t == "read_excel":
        import openpyxl
        path = _sub(cfg.get("file_path", ""), ctx)
        sheet = cfg.get("sheet_name") or None
        wb = openpyxl.load_workbook(path, data_only=True)
        ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            data = []
        else:
            headers = [str(h) if h is not None else f"col{i}" for i, h in enumerate(rows[0])]
            data = [dict(zip(headers, r)) for r in rows[1:]]
            if cfg.get("create_column_vars"):
                for h in headers:
                    ctx["vars"][_slugify_var(h)] = json.dumps([row.get(h) for row in data], ensure_ascii=False, default=str)
        result = json.dumps(data, ensure_ascii=False, default=str)
        _store(result, var, ctx)
        return result

    if t == "write_excel":
        import openpyxl
        dst = _sub(cfg.get("dest_path", ""), ctx)
        data = json.loads(_sub(cfg.get("data_input", "{output}"), ctx))
        if not isinstance(data, list):
            raise Exception("data_input precisa ser uma lista de objetos JSON")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = cfg.get("sheet_name") or "Sheet1"
        if data:
            headers = list(data[0].keys())
            ws.append(headers)
            for row in data:
                ws.append([row.get(h, "") for h in headers])
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        wb.save(dst)
        result = f"Excel salvo em: {dst} ({len(data)} linha(s))"
        _store(result, var, ctx)
        return result

    if t == "write_row":
        import openpyxl
        path = _sub(cfg.get("file_path", ""), ctx)
        sheet = cfg.get("sheet_name") or None
        row_values = json.loads(_sub(cfg.get("data_input", "[]"), ctx))
        if not isinstance(row_values, list):
            raise Exception('data_input precisa ser uma lista JSON de valores (uma linha), ex: ["Nome", 10]')
        row_index = int(cfg.get("row_index", 0) or 0)
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
        if row_index > 0:
            for col_i, val in enumerate(row_values, start=1):
                ws.cell(row=row_index, column=col_i, value=val)
            result = f"Linha {row_index} escrita em: {path}"
        else:
            ws.append(row_values)
            result = f"Linha adicionada ao final de: {path}"
        wb.save(path)
        _store(result, var, ctx)
        return result

    if t == "write_cell":
        import openpyxl
        path = _sub(cfg.get("file_path", ""), ctx)
        sheet = cfg.get("sheet_name") or None
        cell_ref = _sub(cfg.get("cell_ref", ""), ctx).strip().upper()
        if not cell_ref:
            raise Exception("Informe a célula (ex: B3)")
        value = _sub(cfg.get("cell_value", ""), ctx)
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
        ws[cell_ref] = value
        wb.save(path)
        result = f"Célula {cell_ref} = '{value}' em: {path}"
        _store(result, var, ctx)
        return result

    if t == "remove_row":
        import openpyxl
        path = _sub(cfg.get("file_path", ""), ctx)
        sheet = cfg.get("sheet_name") or None
        row_index = int(cfg.get("row_index", 0) or 0)
        if row_index < 1:
            raise Exception("Informe o número da linha a remover (>= 1)")
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
        ws.delete_rows(row_index)
        wb.save(path)
        result = f"Linha {row_index} removida de: {path}"
        _store(result, var, ctx)
        return result

    if t == "remove_cell":
        import openpyxl
        path = _sub(cfg.get("file_path", ""), ctx)
        sheet = cfg.get("sheet_name") or None
        cell_ref = _sub(cfg.get("cell_ref", ""), ctx).strip().upper()
        if not cell_ref:
            raise Exception("Informe a célula (ex: B3)")
        wb = openpyxl.load_workbook(path)
        ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
        ws[cell_ref] = None
        wb.save(path)
        result = f"Célula {cell_ref} limpa em: {path}"
        _store(result, var, ctx)
        return result

    if t == "read_csv":
        import csv as csv_module
        path = _sub(cfg.get("file_path", ""), ctx)
        delim = cfg.get("delimiter", ",") or ","
        with open(path, "r", encoding="utf-8", newline="") as f:
            reader = csv_module.DictReader(f, delimiter=delim)
            data = list(reader)
            headers = reader.fieldnames or []
        if cfg.get("create_column_vars"):
            for h in headers:
                ctx["vars"][_slugify_var(h)] = json.dumps([row.get(h) for row in data], ensure_ascii=False)
        result = json.dumps(data, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "write_csv":
        import csv as csv_module
        dst = _sub(cfg.get("dest_path", ""), ctx)
        delim = cfg.get("delimiter", ",") or ","
        data = json.loads(_sub(cfg.get("data_input", "{output}"), ctx))
        if not isinstance(data, list):
            raise Exception("data_input precisa ser uma lista de objetos JSON")
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        with open(dst, "w", encoding="utf-8", newline="") as f:
            if data:
                writer = csv_module.DictWriter(f, fieldnames=list(data[0].keys()), delimiter=delim)
                writer.writeheader()
                writer.writerows(data)
        result = f"CSV salvo em: {dst} ({len(data)} linha(s))"
        _store(result, var, ctx)
        return result

    if t == "filter_data":
        data = json.loads(_sub(cfg.get("data_input", "{output}"), ctx))
        if not isinstance(data, list):
            raise Exception("data_input precisa ser uma lista de objetos JSON")
        col = cfg.get("sort_key", "") or cfg.get("merge_key", "")
        operator = cfg.get("operator", "contains")
        value = cfg.get("condition_value", "")
        filtered = [row for row in data if isinstance(row, dict) and _eval_condition(str(row.get(col, "")), operator, value)]
        result = json.dumps(filtered, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "merge_data":
        left = json.loads(_sub(cfg.get("data_input", "{output}"), ctx))
        right = json.loads(_sub(cfg.get("data_input2", "[]"), ctx))
        key = cfg.get("merge_key", "")
        if not key:
            raise Exception("Informe a chave de junção (merge_key)")
        right_by_key = {str(r.get(key)): r for r in right if isinstance(r, dict)}
        merged = []
        for row in left:
            match = right_by_key.get(str(row.get(key)), {})
            merged.append({**row, **{k: v for k, v in match.items() if k != key}})
        result = json.dumps(merged, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "dedupe_data":
        data = json.loads(_sub(cfg.get("data_input", "{output}"), ctx))
        key = cfg.get("merge_key", "")
        seen = set()
        deduped = []
        for row in data:
            k = row.get(key) if key and isinstance(row, dict) else json.dumps(row, sort_keys=True, ensure_ascii=False)
            if k in seen:
                continue
            seen.add(k)
            deduped.append(row)
        result = json.dumps(deduped, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "sort_group_data":
        data = json.loads(_sub(cfg.get("data_input", "{output}"), ctx))
        sort_key = cfg.get("sort_key", "")
        desc = bool(cfg.get("sort_desc", False))
        if sort_key:
            data = sorted(data, key=lambda r: (r.get(sort_key) is None, r.get(sort_key)), reverse=desc)
        result = json.dumps(data, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    # ── PDF ─────────────────────────────────────────────────

    if t == "pdf_extract_text":
        import pdfplumber
        path = _sub(cfg.get("source_path", "") or cfg.get("file_path", ""), ctx)
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text_parts.append(page.extract_text() or "")
        result = "\n".join(text_parts)
        _store(result, var, ctx)
        return result

    if t == "pdf_extract_tables":
        import pdfplumber
        path = _sub(cfg.get("source_path", "") or cfg.get("file_path", ""), ctx)
        all_tables = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                for table in page.extract_tables():
                    all_tables.append(table)
        result = json.dumps(all_tables, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "pdf_merge":
        from pypdf import PdfWriter
        raw_list = _sub(cfg.get("list_source", "[]"), ctx)
        paths = json.loads(raw_list)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        writer = PdfWriter()
        for p in paths:
            writer.append(p)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        with open(dst, "wb") as f:
            writer.write(f)
        result = f"PDFs mesclados em: {dst} ({len(paths)} arquivo(s))"
        _store(result, var, ctx)
        return result

    if t == "pdf_split":
        from pypdf import PdfReader, PdfWriter
        path = _sub(cfg.get("source_path", "") or cfg.get("file_path", ""), ctx)
        dst_dir = _sub(cfg.get("dest_path", "") or ".", ctx)
        os.makedirs(dst_dir, exist_ok=True)
        reader = PdfReader(path)
        out_paths = []
        for i, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)
            out_path = os.path.join(dst_dir, f"pagina_{i + 1}.pdf")
            with open(out_path, "wb") as f:
                writer.write(f)
            out_paths.append(out_path)
        result = json.dumps(out_paths, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "pdf_generate":
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas as pdf_canvas
        dst = _sub(cfg.get("dest_path", ""), ctx)
        text = _sub(cfg.get("content", "{output}"), ctx)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        c = pdf_canvas.Canvas(dst, pagesize=A4)
        width, height = A4
        y = height - 50
        for line in text.split("\n"):
            if y < 50:
                c.showPage()
                y = height - 50
            c.drawString(50, y, line[:110])
            y -= 16
        c.save()
        result = f"PDF gerado em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "pdf_fill_form":
        from pypdf import PdfReader, PdfWriter
        path = _sub(cfg.get("source_path", "") or cfg.get("file_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        fields = json.loads(_sub(cfg.get("data_input", "{}"), ctx))
        reader = PdfReader(path)
        writer = PdfWriter()
        writer.append(reader)
        for page in writer.pages:
            writer.update_page_form_field_values(page, fields)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        with open(dst, "wb") as f:
            writer.write(f)
        result = f"Formulário preenchido em: {dst}"
        _store(result, var, ctx)
        return result

    # ── Dados & ETL ─────────────────────────────────────────

    if t == "validate_json_schema":
        import jsonschema
        data = json.loads(_sub(cfg.get("json_input", "{output}"), ctx))
        schema = json.loads(_sub(cfg.get("schema_input", "{}"), ctx))
        try:
            jsonschema.validate(instance=data, schema=schema)
            result = "válido"
        except jsonschema.ValidationError as e:
            result = f"inválido: {e.message}"
        _store(result, var, ctx)
        return result

    if t == "convert_data_format":
        import csv as csv_module
        import io as io_module
        import xmltodict
        import yaml
        raw = _sub(cfg.get("data_input", "{output}"), ctx)
        fmt_from = cfg.get("format_from", "json")
        fmt_to = cfg.get("format_to", "csv")

        if fmt_from == "json":
            data = json.loads(raw)
        elif fmt_from == "yaml":
            data = yaml.safe_load(raw)
        elif fmt_from == "csv":
            data = list(csv_module.DictReader(io_module.StringIO(raw)))
        elif fmt_from == "xml":
            data = xmltodict.parse(raw)
        else:
            raise Exception(f"Formato de origem não suportado: {fmt_from}")

        if fmt_to == "json":
            result = json.dumps(data, ensure_ascii=False, indent=2)
        elif fmt_to == "yaml":
            result = yaml.safe_dump(data, allow_unicode=True)
        elif fmt_to == "csv":
            buf = io_module.StringIO()
            if isinstance(data, list) and data:
                writer = csv_module.DictWriter(buf, fieldnames=list(data[0].keys()))
                writer.writeheader()
                writer.writerows(data)
            result = buf.getvalue()
        elif fmt_to == "xml":
            result = xmltodict.unparse({"root": data}, pretty=True)
        else:
            raise Exception(f"Formato de destino não suportado: {fmt_to}")
        _store(result, var, ctx)
        return result

    if t == "html_extract":
        from bs4 import BeautifulSoup
        html = _sub(cfg.get("text_input", "{output}"), ctx)
        selector = cfg.get("css_selector", "")
        soup = BeautifulSoup(html, "html.parser")
        found = soup.select(selector) if selector else []
        result = json.dumps([el.get_text(strip=True) for el in found], ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "sql_on_data":
        import duckdb
        import pandas as pd
        data = json.loads(_sub(cfg.get("data_input", "{output}"), ctx))
        query = _sub(cfg.get("sql_query", "SELECT * FROM data"), ctx)
        df = pd.DataFrame(data if isinstance(data, list) else [data])
        con = duckdb.connect(":memory:")
        try:
            con.register("data", df)
            rows = con.execute(query).fetchdf()
            result = rows.to_json(orient="records", force_ascii=False)
        finally:
            con.close()
        _store(result, var, ctx)
        return result

    if t == "generate_fake_data":
        from faker import Faker
        fake = Faker("pt_BR")
        kind = cfg.get("fake_type", "name")
        count = int(cfg.get("fake_count", 5) or 5)
        generators = {
            "name": fake.name, "email": fake.email, "cpf": fake.cpf, "cnpj": fake.cnpj,
            "phone": fake.phone_number, "address": fake.address, "company": fake.company,
            "date": lambda: fake.date_between(start_date="-5y", end_date="today").isoformat(),
            "text": lambda: fake.text(max_nb_chars=120),
        }
        gen = generators.get(kind, fake.name)
        result = json.dumps([gen() for _ in range(max(1, min(count, 1000)))], ensure_ascii=False)
        _store(result, var, ctx)
        return result

    # ── Validação BR ────────────────────────────────────────

    if t == "validate_cpf_cnpj":
        from validate_docbr import CPF, CNPJ
        value = re.sub(r"\D", "", _sub(cfg.get("text_input", "{output}"), ctx))
        if len(value) <= 11:
            result = "válido" if CPF().validate(value) else "inválido"
        else:
            result = "válido" if CNPJ().validate(value) else "inválido"
        _store(result, var, ctx)
        return result

    if t == "validate_email":
        from email_validator import validate_email as _validate_email, EmailNotValidError
        addr = _sub(cfg.get("text_input", "{output}"), ctx).strip()
        try:
            _validate_email(addr, check_deliverability=True)
            result = "válido"
        except EmailNotValidError as e:
            result = f"inválido: {e}"
        _store(result, var, ctx)
        return result

    if t == "validate_phone":
        import phonenumbers
        raw = _sub(cfg.get("text_input", "{output}"), ctx)
        region = cfg.get("region", "BR") or "BR"
        try:
            parsed = phonenumbers.parse(raw, region)
            result = "válido" if phonenumbers.is_valid_number(parsed) else "inválido"
        except phonenumbers.NumberParseException as e:
            result = f"inválido: {e}"
        _store(result, var, ctx)
        return result

    if t == "lookup_cep":
        cep = re.sub(r"\D", "", _sub(cfg.get("text_input", "{output}"), ctx))
        async with httpx.AsyncClient(timeout=10) as client:
            resp = await client.get(f"https://viacep.com.br/ws/{cep}/json/")
        data = resp.json()
        if data.get("erro"):
            raise Exception(f"CEP '{cep}' não encontrado")
        result = json.dumps(data, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "format_currency":
        from babel.numbers import format_currency as _format_currency
        raw = _sub(cfg.get("text_input", "{output}"), ctx)
        try:
            value = float(str(raw).replace(",", "."))
        except ValueError:
            raise Exception(f"Valor inválido para formatar: '{raw}'")
        result = _format_currency(value, "BRL", locale="pt_BR")
        _store(result, var, ctx)
        return result

    # ── Segurança & Criptografia ────────────────────────────

    if t == "encrypt_text":
        from cryptography.fernet import Fernet
        import base64 as base64_module
        import hashlib as hashlib_module
        key_raw = _sub(cfg.get("secret_key", ""), ctx)
        if not key_raw:
            raise Exception("Informe uma chave secreta (secret_key)")
        key = base64_module.urlsafe_b64encode(hashlib_module.sha256(key_raw.encode()).digest())
        text = _sub(cfg.get("text_input", "{output}"), ctx)
        result = Fernet(key).encrypt(text.encode()).decode()
        _store(result, var, ctx)
        return result

    if t == "decrypt_text":
        from cryptography.fernet import Fernet, InvalidToken
        import base64 as base64_module
        import hashlib as hashlib_module
        key_raw = _sub(cfg.get("secret_key", ""), ctx)
        if not key_raw:
            raise Exception("Informe uma chave secreta (secret_key)")
        key = base64_module.urlsafe_b64encode(hashlib_module.sha256(key_raw.encode()).digest())
        token = _sub(cfg.get("text_input", "{output}"), ctx)
        try:
            result = Fernet(key).decrypt(token.encode()).decode()
        except InvalidToken:
            raise Exception("Token inválido ou chave secreta incorreta")
        _store(result, var, ctx)
        return result

    if t == "generate_jwt":
        from jose import jwt as jose_jwt
        payload = json.loads(_sub(cfg.get("json_input", "{}"), ctx))
        secret = _sub(cfg.get("secret_key", ""), ctx)
        if not secret:
            raise Exception("Informe uma chave secreta (secret_key)")
        result = jose_jwt.encode(payload, secret, algorithm="HS256")
        _store(result, var, ctx)
        return result

    if t == "verify_jwt":
        from jose import jwt as jose_jwt
        from jose.exceptions import JWTError
        token = _sub(cfg.get("text_input", "{output}"), ctx)
        secret = _sub(cfg.get("secret_key", ""), ctx)
        try:
            payload = jose_jwt.decode(token, secret, algorithms=["HS256"])
            result = json.dumps(payload, ensure_ascii=False)
        except JWTError as e:
            result = f"inválido: {e}"
        _store(result, var, ctx)
        return result

    if t == "hash_password":
        import bcrypt as bcrypt_module
        text = _sub(cfg.get("text_input", "{output}"), ctx)
        result = bcrypt_module.hashpw(text.encode(), bcrypt_module.gensalt()).decode()
        _store(result, var, ctx)
        return result

    if t == "verify_password":
        import bcrypt as bcrypt_module
        plain = _sub(cfg.get("text_input", "{output}"), ctx)
        hashed = _sub(cfg.get("secret_key", ""), ctx)
        try:
            ok = bcrypt_module.checkpw(plain.encode(), hashed.encode())
        except ValueError:
            ok = False
        result = "true" if ok else "false"
        _store(result, var, ctx)
        return result

    if t == "generate_otp":
        import pyotp
        secret = _sub(cfg.get("secret_key", ""), ctx) or pyotp.random_base32()
        code = pyotp.TOTP(secret).now()
        result = json.dumps({"secret": secret, "code": code}, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "verify_otp":
        import pyotp
        secret = _sub(cfg.get("secret_key", ""), ctx)
        code = _sub(cfg.get("text_input", "{output}"), ctx).strip()
        if not secret:
            raise Exception("Informe o secret do OTP (secret_key)")
        ok = pyotp.TOTP(secret).verify(code, valid_window=1)
        result = "true" if ok else "false"
        _store(result, var, ctx)
        return result

    if t == "generate_secure_password":
        length = int(cfg.get("password_length", 16) or 16)
        alphabet = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789!@#$%^&*()-_=+"
        result = "".join(secrets.choice(alphabet) for _ in range(max(4, min(length, 128))))
        _store(result, var, ctx)
        return result

    if t == "check_ssl_cert":
        import ssl as ssl_module
        import socket as socket_module
        from datetime import datetime as _dt
        host = _sub(cfg.get("text_input", "{output}"), ctx).strip()
        ctx_ssl = ssl_module.create_default_context()
        with socket_module.create_connection((host, 443), timeout=10) as sock:
            with ctx_ssl.wrap_socket(sock, server_hostname=host) as ssock:
                cert = ssock.getpeercert()
        expires = _dt.strptime(cert["notAfter"], "%b %d %H:%M:%S %Y %Z")
        days_left = (expires - _dt.utcnow()).days
        result = json.dumps({
            "subject": dict(x[0] for x in cert.get("subject", [])),
            "issuer": dict(x[0] for x in cert.get("issuer", [])),
            "expires": expires.isoformat(), "days_left": days_left,
        }, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "hmac_sign":
        import hmac as hmac_module
        import hashlib as hashlib_module
        message = _sub(cfg.get("text_input", "{output}"), ctx)
        secret = _sub(cfg.get("secret_key", ""), ctx)
        if not secret:
            raise Exception("Informe uma chave secreta (secret_key)")
        result = hmac_module.new(secret.encode(), message.encode(), hashlib_module.sha256).hexdigest()
        _store(result, var, ctx)
        return result

    # ── Comunicação ─────────────────────────────────────────

    if t == "send_telegram":
        bot_token = _sub(cfg.get("secret_key", ""), ctx)
        chat_id = _sub(cfg.get("to", ""), ctx)
        text = _sub(cfg.get("content", "{output}"), ctx)
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.post(
                f"https://api.telegram.org/bot{bot_token}/sendMessage",
                json={"chat_id": chat_id, "text": text},
            )
        if resp.status_code != 200:
            raise Exception(f"Telegram retornou {resp.status_code}: {resp.text}")
        result = f"Mensagem enviada ao chat {chat_id}"
        _store(result, var, ctx)
        return result

    if t == "send_slack":
        webhook_url = _sub(cfg.get("url", ""), ctx)
        text = _sub(cfg.get("content", "{output}"), ctx)
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.post(webhook_url, json={"text": text})
        if resp.status_code != 200:
            raise Exception(f"Slack retornou {resp.status_code}: {resp.text}")
        result = "Mensagem enviada ao Slack"
        _store(result, var, ctx)
        return result

    if t == "send_discord":
        webhook_url = _sub(cfg.get("url", ""), ctx)
        text = _sub(cfg.get("content", "{output}"), ctx)
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.post(webhook_url, json={"content": text})
        if resp.status_code not in (200, 204):
            raise Exception(f"Discord retornou {resp.status_code}: {resp.text}")
        result = "Mensagem enviada ao Discord"
        _store(result, var, ctx)
        return result

    if t in ("send_whatsapp", "send_sms"):
        sid = _sub(cfg.get("api_key", ""), ctx)
        token = _sub(cfg.get("api_secret", ""), ctx)
        from_num = _sub(cfg.get("from_number", ""), ctx)
        to_num = _sub(cfg.get("to", ""), ctx)
        body = _sub(cfg.get("content", "{output}"), ctx)
        prefix = "whatsapp:" if t == "send_whatsapp" else ""
        async with httpx.AsyncClient(timeout=15, auth=(sid, token)) as client:
            resp = await client.post(
                f"https://api.twilio.com/2010-04-01/Accounts/{sid}/Messages.json",
                data={"From": f"{prefix}{from_num}", "To": f"{prefix}{to_num}", "Body": body},
            )
        if resp.status_code not in (200, 201):
            raise Exception(f"Twilio retornou {resp.status_code}: {resp.text}")
        result = f"{'WhatsApp' if t == 'send_whatsapp' else 'SMS'} enviado para {to_num}"
        _store(result, var, ctx)
        return result

    if t == "read_email_imap":
        import imaplib
        import email as email_module
        host = _sub(cfg.get("url", ""), ctx) or "imap.gmail.com"
        user_addr = _sub(cfg.get("to", ""), ctx)
        password = _sub(cfg.get("secret_key", ""), ctx)
        limit = int(cfg.get("fake_count", 5) or 5)
        imap = imaplib.IMAP4_SSL(host)
        try:
            imap.login(user_addr, password)
            imap.select("INBOX")
            _, msg_ids = imap.search(None, "ALL")
            ids = msg_ids[0].split()[-limit:]
            messages = []
            for mid in reversed(ids):
                _, data = imap.fetch(mid, "(RFC822)")
                msg = email_module.message_from_bytes(data[0][1])
                messages.append({"from": msg.get("From", ""), "subject": msg.get("Subject", ""), "date": msg.get("Date", "")})
        finally:
            imap.logout()
        result = json.dumps(messages, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    # ── Notificações ────────────────────────────────────────

    if t == "send_push_notification":
        api_key = _sub(cfg.get("api_key", ""), ctx)
        app_id = _sub(cfg.get("secret_key", ""), ctx)
        message = _sub(cfg.get("content", "{output}"), ctx)
        segment_or_player = _sub(cfg.get("to", ""), ctx)
        payload = {"app_id": app_id, "contents": {"en": message}}
        if segment_or_player:
            payload["include_player_ids"] = [segment_or_player]
        else:
            payload["included_segments"] = ["All"]
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.post(
                "https://onesignal.com/api/v1/notifications",
                headers={"Authorization": f"Basic {api_key}", "Content-Type": "application/json"},
                json=payload,
            )
        if resp.status_code not in (200, 201):
            raise Exception(f"OneSignal retornou {resp.status_code}: {resp.text}")
        result = "Push notification enviada"
        _store(result, var, ctx)
        return result

    if t == "create_incident":
        routing_key = _sub(cfg.get("api_key", ""), ctx)
        summary = _sub(cfg.get("content", "{output}"), ctx)
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.post(
                "https://events.pagerduty.com/v2/enqueue",
                json={
                    "routing_key": routing_key, "event_action": "trigger",
                    "payload": {"summary": summary, "source": "HAC Studio", "severity": "critical"},
                },
            )
        if resp.status_code not in (200, 202):
            raise Exception(f"PagerDuty retornou {resp.status_code}: {resp.text}")
        result = "Incidente criado no PagerDuty"
        _store(result, var, ctx)
        return result

    # ── Pagamentos & Financeiro ─────────────────────────────

    if t == "asaas_create_charge":
        api_key = _sub(cfg.get("api_key", ""), ctx)
        customer_id = _sub(cfg.get("to", ""), ctx)
        value = _sub(cfg.get("text_input", "0"), ctx)
        description = _sub(cfg.get("content", ""), ctx)
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.post(
                "https://api.asaas.com/v3/payments",
                headers={"access_token": api_key, "Content-Type": "application/json"},
                json={"customer": customer_id, "billingType": "PIX", "value": float(value), "description": description},
            )
        if resp.status_code not in (200, 201):
            raise Exception(f"Asaas retornou {resp.status_code}: {resp.text}")
        result = resp.text
        _store(result, var, ctx)
        return result

    if t == "asaas_check_payment":
        api_key = _sub(cfg.get("api_key", ""), ctx)
        payment_id = _sub(cfg.get("text_input", "{output}"), ctx)
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                f"https://api.asaas.com/v3/payments/{payment_id}",
                headers={"access_token": api_key},
            )
        if resp.status_code != 200:
            raise Exception(f"Asaas retornou {resp.status_code}: {resp.text}")
        result = resp.text
        _store(result, var, ctx)
        return result

    if t == "generate_pix_qr":
        key = _sub(cfg.get("pix_key", ""), ctx)
        name = _sub(cfg.get("pix_merchant_name", ""), ctx)
        city = _sub(cfg.get("pix_merchant_city", ""), ctx)
        amount = _sub(cfg.get("text_input", ""), ctx)
        if not key:
            raise Exception("Informe a chave Pix (pix_key)")
        result = _build_pix_payload(key, name, city, amount)
        _store(result, var, ctx)
        return result

    if t == "get_currency_rate":
        pair = _sub(cfg.get("text_input", "USD-BRL"), ctx).strip()
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(f"https://economia.awesomeapi.com.br/last/{pair}")
        data = resp.json()
        result = json.dumps(data, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "get_crypto_price":
        coin = _sub(cfg.get("text_input", "bitcoin"), ctx).strip().lower()
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                "https://api.coingecko.com/api/v3/simple/price",
                params={"ids": coin, "vs_currencies": "usd,brl"},
            )
        data = resp.json()
        if coin not in data:
            raise Exception(f"Moeda '{coin}' não encontrada")
        result = json.dumps(data[coin], ensure_ascii=False)
        _store(result, var, ctx)
        return result

    # ── APIs Externas Úteis ─────────────────────────────────

    if t == "get_weather":
        city = _sub(cfg.get("text_input", "{output}"), ctx)
        api_key = _sub(cfg.get("api_key", ""), ctx)
        if not api_key:
            raise Exception("Informe a API key do OpenWeatherMap (api_key)")
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(
                "https://api.openweathermap.org/data/2.5/weather",
                params={"q": city, "appid": api_key, "units": "metric", "lang": "pt_br"},
            )
        if resp.status_code != 200:
            raise Exception(f"OpenWeatherMap retornou {resp.status_code}: {resp.text}")
        result = resp.text
        _store(result, var, ctx)
        return result

    if t == "geocode_address":
        address = _sub(cfg.get("text_input", "{output}"), ctx)
        async with httpx.AsyncClient(timeout=15, headers={"User-Agent": "HAC-Studio/1.0"}) as client:
            resp = await client.get(
                "https://nominatim.openstreetmap.org/search",
                params={"q": address, "format": "json", "limit": 1},
            )
        data = resp.json()
        if not data:
            raise Exception(f"Endereço não encontrado: '{address}'")
        result = json.dumps({"lat": data[0]["lat"], "lon": data[0]["lon"], "display_name": data[0]["display_name"]}, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "calculate_distance":
        import math
        lat1, lon1 = [float(x) for x in _sub(cfg.get("coord_from", ""), ctx).split(",")]
        lat2, lon2 = [float(x) for x in _sub(cfg.get("coord_to", ""), ctx).split(",")]
        r = 6371.0
        phi1, phi2 = math.radians(lat1), math.radians(lat2)
        dphi = math.radians(lat2 - lat1)
        dlambda = math.radians(lon2 - lon1)
        a = math.sin(dphi / 2) ** 2 + math.cos(phi1) * math.cos(phi2) * math.sin(dlambda / 2) ** 2
        km = r * 2 * math.atan2(math.sqrt(a), math.sqrt(1 - a))
        result = str(round(km, 3))
        _store(result, var, ctx)
        return result

    if t == "shorten_url":
        long_url = _sub(cfg.get("text_input", "{output}"), ctx)
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get("http://tinyurl.com/api-create.php", params={"url": long_url})
        if resp.status_code != 200:
            raise Exception(f"TinyURL retornou {resp.status_code}: {resp.text}")
        result = resp.text.strip()
        _store(result, var, ctx)
        return result

    if t == "lookup_cnpj":
        cnpj = re.sub(r"\D", "", _sub(cfg.get("text_input", "{output}"), ctx))
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(f"https://brasilapi.com.br/api/cnpj/v1/{cnpj}")
        if resp.status_code != 200:
            raise Exception(f"CNPJ '{cnpj}' não encontrado ou inválido")
        result = resp.text
        _store(result, var, ctx)
        return result

    if t == "translate_text":
        text = _sub(cfg.get("text_input", "{output}"), ctx)
        target_lang = (cfg.get("region", "EN") or "EN").upper()
        api_key = _sub(cfg.get("api_key", ""), ctx)
        if not api_key:
            raise Exception("Informe a API key do DeepL (api_key)")
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.post(
                "https://api-free.deepl.com/v2/translate",
                headers={"Authorization": f"DeepL-Auth-Key {api_key}"},
                data={"text": text, "target_lang": target_lang},
            )
        if resp.status_code != 200:
            raise Exception(f"DeepL retornou {resp.status_code}: {resp.text}")
        result = resp.json()["translations"][0]["text"]
        _store(result, var, ctx)
        return result

    if t == "get_holidays":
        year = _sub(cfg.get("text_input", "") or str(datetime.utcnow().year), ctx).strip()
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.get(f"https://brasilapi.com.br/api/feriados/v1/{year}")
        if resp.status_code != 200:
            raise Exception(f"Ano inválido ou API indisponível: {resp.status_code}")
        result = resp.text
        _store(result, var, ctx)
        return result

    # ── Web & Scraping avançado ─────────────────────────────

    if t == "download_file":
        url = _sub(cfg.get("url", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        async with httpx.AsyncClient(timeout=60, follow_redirects=True) as client:
            async with client.stream("GET", url) as resp:
                if resp.status_code != 200:
                    raise Exception(f"Download falhou: {resp.status_code}")
                with open(dst, "wb") as f:
                    async for chunk in resp.aiter_bytes():
                        f.write(chunk)
        result = f"Baixado em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "upload_file":
        url = _sub(cfg.get("url", ""), ctx)
        path = _sub(cfg.get("source_path", ""), ctx)
        headers = cfg.get("headers") or {}
        with open(path, "rb") as f:
            async with httpx.AsyncClient(timeout=60) as client:
                resp = await client.post(url, files={"file": (os.path.basename(path), f)}, headers=headers)
        result = f"[{resp.status_code}] {resp.text[:2000]}"
        _store(result, var, ctx)
        return result

    if t == "scrape_html_table":
        import pandas as pd
        source = _sub(cfg.get("text_input", "{output}"), ctx)
        if source.startswith("http://") or source.startswith("https://"):
            async with httpx.AsyncClient(timeout=20, headers={"User-Agent": "Mozilla/5.0"}) as client:
                resp = await client.get(source)
            html = resp.text
        else:
            html = source
        tables = pd.read_html(io.StringIO(html))
        result = json.dumps([t_df.to_dict(orient="records") for t_df in tables], ensure_ascii=False, default=str)
        _store(result, var, ctx)
        return result

    if t == "read_rss_feed":
        import feedparser
        url = _sub(cfg.get("url", ""), ctx)
        feed = feedparser.parse(url)
        items = [{"title": e.get("title", ""), "link": e.get("link", ""), "published": e.get("published", "")} for e in feed.entries]
        result = json.dumps(items, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "http_request_retry":
        method = cfg.get("method", "GET").upper()
        url = _sub(cfg.get("url", ""), ctx)
        headers = cfg.get("headers") or {}
        body = _sub(cfg.get("body", ""), ctx)
        max_tries = int(cfg.get("max_iterations", 3) or 3)
        last_exc = None
        for attempt in range(max_tries):
            try:
                async with httpx.AsyncClient(timeout=30) as client:
                    resp = await client.request(method, url, headers=headers, content=body or None)
                if resp.status_code < 500:
                    result = f"[{resp.status_code}] {resp.text[:3000]}"
                    _store(resp.text[:5000], var, ctx)
                    return result
                last_exc = Exception(f"HTTP {resp.status_code}")
            except Exception as e:
                last_exc = e
            if attempt < max_tries - 1:
                await asyncio.sleep(min(2 ** attempt, 30))
        raise Exception(f"Falhou após {max_tries} tentativas: {last_exc}")

    # ── Texto & NLP ─────────────────────────────────────────

    if t == "detect_language":
        from langdetect import detect
        text = _sub(cfg.get("text_input", "{output}"), ctx)
        try:
            result = detect(text)
        except Exception:
            result = "desconhecido"
        _store(result, var, ctx)
        return result

    if t == "count_tokens":
        import tiktoken
        text = _sub(cfg.get("text_input", "{output}"), ctx)
        enc = tiktoken.get_encoding("cl100k_base")
        result = str(len(enc.encode(text)))
        _store(result, var, ctx)
        return result

    # ── IA & ML extra ───────────────────────────────────────

    if t == "generate_embedding":
        from openai import AsyncOpenAI
        text = _sub(cfg.get("text_input", "{output}"), ctx)
        api_key = _sub(cfg.get("api_key", ""), ctx)
        if not api_key:
            raise Exception("Informe a API key da OpenAI (api_key)")
        client = AsyncOpenAI(api_key=api_key)
        resp = await client.embeddings.create(model="text-embedding-3-small", input=text)
        result = json.dumps(resp.data[0].embedding, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "semantic_search":
        from openai import AsyncOpenAI
        import numpy as np
        query = _sub(cfg.get("text_input", "{output}"), ctx)
        api_key = _sub(cfg.get("api_key", ""), ctx)
        candidates = json.loads(_sub(cfg.get("json_input", "[]"), ctx))
        if not api_key:
            raise Exception("Informe a API key da OpenAI (api_key)")
        if not candidates:
            raise Exception("Informe candidatos em json_input: [{\"text\":..,\"embedding\":[..]}]")
        client = AsyncOpenAI(api_key=api_key)
        q_resp = await client.embeddings.create(model="text-embedding-3-small", input=query)
        q_vec = np.array(q_resp.data[0].embedding)
        best, best_score = None, -1.0
        for cand in candidates:
            vec = np.array(cand["embedding"])
            score = float(np.dot(q_vec, vec) / (np.linalg.norm(q_vec) * np.linalg.norm(vec) + 1e-9))
            if score > best_score:
                best, best_score = cand, score
        result = json.dumps({"text": best.get("text"), "score": round(best_score, 4)}, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "moderate_content":
        from openai import AsyncOpenAI
        text = _sub(cfg.get("text_input", "{output}"), ctx)
        api_key = _sub(cfg.get("api_key", ""), ctx)
        if not api_key:
            raise Exception("Informe a API key da OpenAI (api_key)")
        client = AsyncOpenAI(api_key=api_key)
        resp = await client.moderations.create(input=text)
        r0 = resp.results[0]
        result = json.dumps({"flagged": r0.flagged, "categories": {k: v for k, v in r0.categories.model_dump().items() if v}}, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "compare_texts":
        import difflib
        text_a = _sub(cfg.get("text_input", "{output}"), ctx)
        text_b = _sub(cfg.get("data_input2", ""), ctx)
        ratio = difflib.SequenceMatcher(None, text_a, text_b).ratio()
        result = str(round(ratio, 4))
        _store(result, var, ctx)
        return result

    # ── Sistema & DevOps ─────────────────────────────────────

    if t == "system_stats":
        import psutil
        stats = {
            "cpu_percent": psutil.cpu_percent(interval=0.3),
            "memory_percent": psutil.virtual_memory().percent,
            "disk_percent": psutil.disk_usage(os.path.abspath(os.sep)).percent,
        }
        result = json.dumps(stats, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "list_processes":
        import psutil
        limit = int(cfg.get("fake_count", 20) or 20)
        procs = []
        for p in psutil.process_iter(["pid", "name", "cpu_percent", "memory_percent"]):
            try:
                procs.append(p.info)
            except (psutil.NoSuchProcess, psutil.AccessDenied):
                continue
        procs = sorted(procs, key=lambda p: p.get("memory_percent") or 0, reverse=True)[:max(1, min(limit, 200))]
        result = json.dumps(procs, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "check_port_open":
        import socket as socket_module
        raw = _sub(cfg.get("text_input", "{output}"), ctx)
        host, _, port_s = raw.rpartition(":")
        if not host:
            raise Exception("Informe host:porta, ex: exemplo.com:443")
        port = int(port_s)
        sock = socket_module.socket(socket_module.AF_INET, socket_module.SOCK_STREAM)
        sock.settimeout(5)
        try:
            open_ = sock.connect_ex((host, port)) == 0
        finally:
            sock.close()
        result = "aberta" if open_ else "fechada"
        _store(result, var, ctx)
        return result

    if t == "dns_lookup":
        import dns.resolver
        domain = _sub(cfg.get("text_input", "{output}"), ctx)
        record_type = (cfg.get("operation", "A") or "A").upper()
        answers = dns.resolver.resolve(domain, record_type)
        result = json.dumps([str(a) for a in answers], ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "whois_lookup":
        import whois as whois_module
        domain = _sub(cfg.get("text_input", "{output}"), ctx)
        w = whois_module.whois(domain)
        result = json.dumps({k: str(v) for k, v in dict(w).items()}, ensure_ascii=False, default=str)
        _store(result, var, ctx)
        return result

    if t == "ssh_execute":
        import paramiko
        raw_host = _sub(cfg.get("url", ""), ctx)
        host, _, port_s = raw_host.rpartition(":")
        if not host:
            host, port = raw_host, 22
        else:
            try:
                port = int(port_s)
            except ValueError:
                host, port = raw_host, 22
        username = _sub(cfg.get("to", ""), ctx)
        password = _sub(cfg.get("secret_key", ""), ctx)
        command = _sub(cfg.get("command", ""), ctx)
        client = paramiko.SSHClient()
        client.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        try:
            client.connect(host, port=port, username=username, password=password, timeout=10)
            _, stdout, stderr = client.exec_command(command, timeout=30)
            result = stdout.read().decode("utf-8", errors="replace") + stderr.read().decode("utf-8", errors="replace")
        finally:
            client.close()
        _store(result, var, ctx)
        return result

    if t == "read_env_var":
        name = _sub(cfg.get("text_input", "{output}"), ctx).strip()
        result = os.environ.get(name, "")
        _store(result, var, ctx)
        return result

    if t == "check_url_uptime":
        url = _sub(cfg.get("url", ""), ctx)
        start = time.time()
        try:
            async with httpx.AsyncClient(timeout=15) as client:
                resp = await client.get(url)
            up = resp.status_code < 500
            status_code = resp.status_code
        except Exception:
            up, status_code = False, 0
        latency_ms = int((time.time() - start) * 1000)
        result = json.dumps({"up": up, "status_code": status_code, "latency_ms": latency_ms}, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    # ── Banco de Dados & Fila ────────────────────────────────

    if t == "redis_get":
        import redis.asyncio as redis_module
        conn = redis_module.from_url(_sub(cfg.get("url", ""), ctx), decode_responses=True)
        try:
            result = await conn.get(_sub(cfg.get("text_input", ""), ctx)) or ""
        finally:
            await conn.aclose()
        _store(result, var, ctx)
        return result

    if t == "redis_set":
        import redis.asyncio as redis_module
        conn = redis_module.from_url(_sub(cfg.get("url", ""), ctx), decode_responses=True)
        try:
            await conn.set(_sub(cfg.get("text_input", ""), ctx), _sub(cfg.get("content", "{output}"), ctx))
        finally:
            await conn.aclose()
        result = "salvo"
        _store(result, var, ctx)
        return result

    if t == "queue_push":
        import redis.asyncio as redis_module
        conn = redis_module.from_url(_sub(cfg.get("url", ""), ctx), decode_responses=True)
        try:
            await conn.rpush(_sub(cfg.get("text_input", ""), ctx), _sub(cfg.get("content", "{output}"), ctx))
        finally:
            await conn.aclose()
        result = "enfileirado"
        _store(result, var, ctx)
        return result

    if t == "queue_pop":
        import redis.asyncio as redis_module
        conn = redis_module.from_url(_sub(cfg.get("url", ""), ctx), decode_responses=True)
        try:
            result = await conn.lpop(_sub(cfg.get("text_input", ""), ctx)) or ""
        finally:
            await conn.aclose()
        _store(result, var, ctx)
        return result

    if t == "sql_query_external":
        import asyncpg
        dsn = _sub(cfg.get("url", ""), ctx)
        query = _sub(cfg.get("sql_query", "SELECT 1"), ctx)
        conn = await asyncpg.connect(dsn)
        try:
            rows = await conn.fetch(query)
            result = json.dumps([dict(r) for r in rows], ensure_ascii=False, default=str)
        finally:
            await conn.close()
        _store(result, var, ctx)
        return result

    # ── Templates & Documentos ───────────────────────────────

    if t == "render_template":
        from jinja2 import Template
        template_str = cfg.get("content", "") or ""
        variables = json.loads(_sub(cfg.get("json_input", "{}"), ctx))
        result = Template(template_str).render(**variables, output=ctx.get("output", ""), input=ctx.get("input", ""))
        _store(result, var, ctx)
        return result

    if t == "generate_word_doc":
        import docx
        dst = _sub(cfg.get("dest_path", ""), ctx)
        text = _sub(cfg.get("content", "{output}"), ctx)
        doc = docx.Document()
        for line in text.split("\n"):
            doc.add_paragraph(line)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        doc.save(dst)
        result = f"Documento salvo em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "generate_pptx":
        from pptx import Presentation
        dst = _sub(cfg.get("dest_path", ""), ctx)
        slides_data = json.loads(_sub(cfg.get("data_input", "[]"), ctx))
        prs = Presentation()
        layout = prs.slide_layouts[1]
        for slide_info in slides_data:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_info.get("title", "")
            slide.placeholders[1].text = slide_info.get("content", "")
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        prs.save(dst)
        result = f"Apresentação salva em: {dst} ({len(slides_data)} slide(s))"
        _store(result, var, ctx)
        return result

    # ── Imagens ─────────────────────────────────────────────

    if t == "resize_image":
        from PIL import Image
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        width = int(cfg.get("width", 0) or 0)
        height = int(cfg.get("height", 0) or 0)
        img = Image.open(src)
        if width and not height:
            height = int(img.height * (width / img.width))
        elif height and not width:
            width = int(img.width * (height / img.height))
        img = img.resize((width or img.width, height or img.height))
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        img.save(dst)
        result = f"Redimensionado para {width}x{height}: {dst}"
        _store(result, var, ctx)
        return result

    if t == "convert_image_format":
        from PIL import Image
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        img = Image.open(src).convert("RGB")
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        img.save(dst)
        result = f"Convertido para: {dst}"
        _store(result, var, ctx)
        return result

    if t == "add_watermark":
        from PIL import Image, ImageDraw
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        text = _sub(cfg.get("text", ""), ctx)
        img = Image.open(src).convert("RGBA")
        overlay = Image.new("RGBA", img.size, (255, 255, 255, 0))
        draw = ImageDraw.Draw(overlay)
        margin = 10
        draw.text((margin, img.height - 30 - margin), text, fill=(255, 255, 255, 160))
        combined = Image.alpha_composite(img, overlay).convert("RGB")
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        combined.save(dst)
        result = f"Marca d'água aplicada: {dst}"
        _store(result, var, ctx)
        return result

    if t == "generate_thumbnail":
        from PIL import Image
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        size = int(cfg.get("width", 200) or 200)
        img = Image.open(src)
        img.thumbnail((size, size))
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        img.save(dst)
        result = f"Thumbnail salva em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "generate_qrcode":
        import qrcode as qrcode_module
        content = _sub(cfg.get("text_input", "{output}"), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        img = qrcode_module.make(content)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        img.save(dst)
        result = f"QR Code salvo em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "read_qrcode":
        from PIL import Image
        from pyzbar.pyzbar import decode as qr_decode
        src = _sub(cfg.get("source_path", ""), ctx)
        decoded = qr_decode(Image.open(src))
        result = json.dumps([d.data.decode("utf-8", errors="replace") for d in decoded], ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "compare_images":
        import imagehash
        from PIL import Image
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        hash_a = imagehash.average_hash(Image.open(src))
        hash_b = imagehash.average_hash(Image.open(dst))
        diff_bits = hash_a - hash_b
        similarity = round(1 - (diff_bits / len(hash_a.hash) ** 2), 4)
        result = json.dumps({"diff_bits": diff_bits, "similarity": similarity}, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    if t == "generate_ai_image":
        from openai import AsyncOpenAI
        prompt = _sub(cfg.get("text_input", "{output}"), ctx)
        api_key = _sub(cfg.get("api_key", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        if not api_key:
            raise Exception("Informe a API key da OpenAI (api_key)")
        client = AsyncOpenAI(api_key=api_key)
        resp = await client.images.generate(model="dall-e-3", prompt=prompt, n=1, size="1024x1024")
        image_url = resp.data[0].url
        async with httpx.AsyncClient(timeout=60) as http_client:
            img_resp = await http_client.get(image_url)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        with open(dst, "wb") as f:
            f.write(img_resp.content)
        result = f"Imagem gerada em: {dst}"
        _store(result, var, ctx)
        return result

    # ── Áudio & Vídeo (requer ffmpeg instalado no host do agente) ──

    if t == "transcode_media":
        import ffmpeg as ffmpeg_module
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        ffmpeg_module.input(src).output(dst).overwrite_output().run(quiet=True)
        result = f"Transcodificado em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "extract_audio":
        import ffmpeg as ffmpeg_module
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        ffmpeg_module.input(src).output(dst, vn=None).overwrite_output().run(quiet=True)
        result = f"Áudio extraído em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "trim_media":
        import ffmpeg as ffmpeg_module
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        start = float(cfg.get("seconds", 0) or 0)
        duration = float(cfg.get("seconds_max", 0) or 0)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        stream = ffmpeg_module.input(src, ss=start)
        if duration:
            stream = stream.output(dst, t=duration)
        else:
            stream = stream.output(dst)
        stream.overwrite_output().run(quiet=True)
        result = f"Recorte salvo em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "extract_video_frame":
        import ffmpeg as ffmpeg_module
        src = _sub(cfg.get("source_path", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        timestamp = float(cfg.get("seconds", 0) or 0)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        ffmpeg_module.input(src, ss=timestamp).output(dst, vframes=1).overwrite_output().run(quiet=True)
        result = f"Frame salvo em: {dst}"
        _store(result, var, ctx)
        return result

    if t == "transcribe_audio":
        from openai import AsyncOpenAI
        src = _sub(cfg.get("source_path", ""), ctx)
        api_key = _sub(cfg.get("api_key", ""), ctx)
        if not api_key:
            raise Exception("Informe a API key da OpenAI (api_key)")
        client = AsyncOpenAI(api_key=api_key)
        with open(src, "rb") as f:
            transcript = await client.audio.transcriptions.create(model="whisper-1", file=f)
        result = transcript.text
        _store(result, var, ctx)
        return result

    if t == "text_to_speech":
        from openai import AsyncOpenAI
        text = _sub(cfg.get("text_input", "{output}"), ctx)
        api_key = _sub(cfg.get("api_key", ""), ctx)
        dst = _sub(cfg.get("dest_path", ""), ctx)
        if not api_key:
            raise Exception("Informe a API key da OpenAI (api_key)")
        client = AsyncOpenAI(api_key=api_key)
        resp = await client.audio.speech.create(model="tts-1", voice="alloy", input=text)
        if os.path.dirname(dst):
            os.makedirs(os.path.dirname(dst), exist_ok=True)
        resp.write_to_file(dst)
        result = f"Áudio gerado em: {dst}"
        _store(result, var, ctx)
        return result

    # ── OCR & Visão (requer tesseract/poppler no host do agente) ───

    if t == "ocr_image":
        _ensure_native_binary("Tesseract OCR", _tesseract_installed, "UB-Mannheim.TesseractOCR", "tesseract-ocr", "tesseract")
        import pytesseract
        from PIL import Image
        _tess_bin = _find_tesseract_binary()
        if _tess_bin:
            pytesseract.pytesseract.tesseract_cmd = _tess_bin
        src = _sub(cfg.get("source_path", ""), ctx)
        result = pytesseract.image_to_string(Image.open(src), lang="por+eng")
        _store(result, var, ctx)
        return result

    if t == "ocr_pdf_scanned":
        _ensure_native_binary("Tesseract OCR", _tesseract_installed, "UB-Mannheim.TesseractOCR", "tesseract-ocr", "tesseract")
        _ensure_native_binary("Poppler (pdftoppm)", _poppler_installed, "oschwartz10612.Poppler", "poppler-utils", "poppler")
        import pytesseract
        from pdf2image import convert_from_path
        _tess_bin = _find_tesseract_binary()
        if _tess_bin:
            pytesseract.pytesseract.tesseract_cmd = _tess_bin
        _poppler_dir = _find_poppler_bin_dir()
        src = _sub(cfg.get("source_path", ""), ctx)
        pages = convert_from_path(src, poppler_path=_poppler_dir)
        result = "\n\n".join(pytesseract.image_to_string(p, lang="por+eng") for p in pages)
        _store(result, var, ctx)
        return result

    if t == "detect_face_object":
        import cv2
        src = _sub(cfg.get("source_path", ""), ctx)
        cascade = cv2.CascadeClassifier(cv2.data.haarcascades + "haarcascade_frontalface_default.xml")
        img = cv2.imread(src)
        if img is None:
            raise Exception(f"Não foi possível abrir a imagem: {src}")
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        faces = cascade.detectMultiScale(gray, scaleFactor=1.1, minNeighbors=5)
        boxes = [{"x": int(x), "y": int(y), "w": int(w), "h": int(h)} for (x, y, w, h) in faces]
        result = json.dumps({"count": len(boxes), "boxes": boxes}, ensure_ascii=False)
        _store(result, var, ctx)
        return result

    # ── HTTP & Internet ────────────────────────────────────

    if t == "http_request":
        method = cfg.get("method", "GET").upper()
        url = _sub(cfg.get("url", ""), ctx)
        headers = cfg.get("headers") or {}
        body = _sub(cfg.get("body", ""), ctx)

        async with httpx.AsyncClient(timeout=30) as client:
            if method == "GET":
                resp = await client.get(url, headers=headers)
            elif method in ("POST", "PUT", "PATCH"):
                ct = headers.get("Content-Type", "application/json")
                if "json" in ct:
                    try:
                        resp = await client.request(method, url, json=json.loads(body) if body else {}, headers=headers)
                    except json.JSONDecodeError:
                        resp = await client.request(method, url, content=body, headers=headers)
                else:
                    resp = await client.request(method, url, content=body, headers=headers)
            elif method == "DELETE":
                resp = await client.delete(url, headers=headers)
            else:
                resp = await client.request(method, url, content=body, headers=headers)

        result = resp.text
        _store(result[:5000], var, ctx)
        return f"[{resp.status_code}] {result[:3000]}"

    if t == "parse_json":
        raw = _sub(cfg.get("json_input", "{output}"), ctx)
        try:
            data = json.loads(raw)
        except Exception as e:
            raise Exception(f"JSON inválido: {e}")
        key_path = cfg.get("key_path", "")
        if key_path:
            for key in key_path.split("."):
                if isinstance(data, list):
                    try: data = data[int(key)]
                    except: raise Exception(f"Índice inválido '{key}'")
                elif isinstance(data, dict):
                    data = data.get(key)
                    if data is None:
                        raise Exception(f"Chave '{key}' não encontrada")
        result = json.dumps(data, ensure_ascii=False) if isinstance(data, (dict, list)) else str(data)
        _store(result, var, ctx)
        return result

    # ── Email ──────────────────────────────────────────────

    if t == "send_email":
        to_addr = _sub(cfg.get("to", ""), ctx)
        subject = _sub(cfg.get("subject", ""), ctx)
        body_content = _sub(cfg.get("email_body", ""), ctx)
        is_html = cfg.get("is_html", False)

        if not to_addr:
            raise Exception("Destinatário (to) não definido")

        payload = {
            "sender": {
                "name": getattr(settings, "brevo_sender_name", "HAC Studio"),
                "email": getattr(settings, "brevo_sender_email", None) or "no-reply@hacplatform.com",
            },
            "to": [{"email": to_addr}],
            "subject": subject,
        }
        if is_html:
            payload["htmlContent"] = body_content
        else:
            payload["textContent"] = body_content

        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.post(
                "https://api.brevo.com/v3/smtp/email",
                headers={"api-key": settings.brevo_api_key, "Content-Type": "application/json"},
                json=payload,
            )
        if resp.status_code not in (200, 201):
            raise Exception(f"Brevo retornou {resp.status_code}: {resp.text}")
        result = f"Email enviado para {to_addr}"
        _store(result, var, ctx)
        return result

    # ── Sistema ────────────────────────────────────────────

    if t == "run_command":
        command = _sub(cfg.get("command", ""), ctx)
        run_env = os.environ.copy()
        run_env["PYTHONIOENCODING"] = "utf-8"
        if sys.platform == "win32":
            # No Windows o cmd.exe escreve na codepage ANSI ativa (ex: cp1252) por
            # padrão — decodificar como utf-8 sem isso só troca um tipo de acento
            # errado por outro. "chcp 65001" muda a codepage do subshell pra UTF-8
            # antes de rodar o comando de verdade.
            command = f"chcp 65001>nul & {command}"
        proc = subprocess.run(command, shell=True, capture_output=True, text=True,
                               encoding="utf-8", errors="replace", timeout=30, env=run_env)
        result = proc.stdout + proc.stderr
        _store(result, var, ctx)
        return result

    if t == "run_python":
        code = _sub(cfg.get("code", ""), ctx)
        old_stdout = sys.stdout
        sys.stdout = buffer = io.StringIO()
        exec_globals = dict(__builtins__=__builtins__)
        # ctx['vars'] só guarda strings (usadas em substituição de texto {varname} em
        # outros steps), mas aqui os valores viram variáveis Python de verdade — então
        # "0"/"3" (índice de loop/foreach) e listas/objetos JSON (ex: coluna de planilha
        # via 'criar variáveis por coluna') ficam number/list/dict de verdade, não string,
        # senão `minha_lista[item_index]` quebra com "indices must be integers".
        for _k, _v in ctx.get("vars", {}).items():
            try:
                exec_globals[_k] = json.loads(_v)
            except (TypeError, ValueError):
                exec_globals[_k] = _v
        exec_globals["input_data"] = ctx.get("input", "")
        exec_globals["output"] = ctx.get("output", "")
        try:
            exec(code, exec_globals)  # noqa: S102
        finally:
            sys.stdout = old_stdout
        captured = buffer.getvalue()
        # If variable_name specified, try to read that variable from exec_globals
        if var and var in exec_globals:
            result = str(exec_globals[var])
        else:
            result = captured
        _store(result, var, ctx)
        return result

    # ── Inteligência Artificial ────────────────────────────

    if t == "call_ai_agent":
        agent = await ai_agents_col.find_one({"_id": cfg.get("agent_id", "")})
        if not agent:
            raise Exception(f"Agente IA '{cfg.get('agent_id')}' não encontrado")
        template = cfg.get("input_template", "{output}") or "{output}"
        step_input = _sub(template, ctx)
        output_text, _ = await call_ai(agent, step_input)
        _store(output_text, var, ctx)
        return output_text

    if t == "call_pipeline":
        pipeline = await pipelines_col.find_one({"_id": cfg.get("pipeline_id", "")})
        if not pipeline:
            raise Exception(f"Pipeline '{cfg.get('pipeline_id')}' não encontrada")
        template = cfg.get("input_template", "{output}") or "{output}"
        step_input = _sub(template, ctx)
        pipe_output = step_input
        for p_step in pipeline.get("steps", []):
            agent = await ai_agents_col.find_one({"_id": p_step["ai_agent_id"]})
            if not agent:
                raise Exception(f"Agente IA '{p_step['ai_agent_id']}' não encontrado")
            p_tmpl = p_step.get("input_template") or "{output}"
            p_in = p_tmpl.replace("{input}", step_input).replace("{output}", pipe_output)
            pipe_output, _ = await call_ai(agent, p_in)
        _store(pipe_output, var, ctx)
        return pipe_output

    # ── Dados ──────────────────────────────────────────────

    if t == "text_transform":
        text = _sub(cfg.get("text_input", "{output}"), ctx)
        operation = cfg.get("operation", "upper")
        if operation == "upper":      result = text.upper()
        elif operation == "lower":    result = text.lower()
        elif operation == "strip":    result = text.strip()
        elif operation == "replace":
            result = text.replace(cfg.get("search", ""), cfg.get("replace_with", ""))
        elif operation == "count_chars":  result = str(len(text))
        elif operation == "count_words":  result = str(len(text.split()))
        elif operation == "split":
            delim = cfg.get("search", "\n")
            result = json.dumps(text.split(delim), ensure_ascii=False)
        elif operation == "regex":
            pattern = cfg.get("search", "")
            matches = re.findall(pattern, text)
            result = json.dumps(matches, ensure_ascii=False)
        elif operation == "base64_encode":
            import base64
            result = base64.b64encode(text.encode()).decode()
        elif operation == "base64_decode":
            import base64
            result = base64.b64decode(text.encode()).decode()
        elif operation == "remove_accents":
            from unidecode import unidecode
            result = unidecode(text)
        elif operation == "slugify":
            from slugify import slugify
            result = slugify(text)
        else:
            result = text
        _store(result, var, ctx)
        return result

    # ── Navegador ──────────────────────────────────────────

    if t == "browser":
        actions  = cfg.get("browser_actions") or []
        engine   = cfg.get("browser_engine", "playwright")
        headless = cfg.get("browser_headless", True)
        if not actions:
            result = "Nenhuma ação de navegador configurada."
            _store(result, var, ctx)
            return result
        agent_id = ctx.get("agent_id", "")
        user_id  = ctx.get("user_id", "")
        if not agent_id:
            raise Exception(
                "Step Browser requer um agente selecionado. "
                "Escolha um agente no seletor ao lado do botão Executar antes de salvar."
            )
        result_str, var_updates = await _exec_browser_via_agent(actions, engine, headless, ctx, agent_id, user_id)
        ctx["vars"].update(var_updates)
        _store(result_str, var, ctx)
        return result_str

    # ── Navegador (sessão persistente) ─────────────────────

    if t in ("browser_open", "browser_click", "browser_type", "browser_extract",
             "browser_wait", "browser_screenshot", "browser_close",
             "browser_captcha_detect", "browser_captcha_wait", "browser_captcha_solve_image"):
        session_name = _sub(cfg.get("session_name", ""), ctx).strip()
        if not session_name:
            raise Exception("Informe um nome para a sessão no campo 'Nome da sessão'.")

        agent_id = ctx.get("agent_id", "")
        user_id  = ctx.get("user_id", "")
        if not agent_id:
            raise Exception(
                "Esta etapa requer um agente selecionado. "
                "Escolha um agente no seletor ao lado do botão Executar antes de salvar."
            )

        sessions = ctx.setdefault("sessions", {})

        if t == "browser_open":
            if session_name in sessions:
                result = f"Sessão '{session_name}' já está aberta — reaproveitando."
                _store(result, var, ctx)
                return result
            headless = cfg.get("browser_headless", True)
            engine = cfg.get("browser_engine", "playwright")
            if engine not in ("playwright", "selenium"):
                engine = "playwright"
            profile_name = _sub(cfg.get("browser_profile", ""), ctx).strip()
            script = _gen_session_open_script(session_name, cfg.get("target", ""), headless, ctx, engine=engine,
                                              profile_name=profile_name)
            raw = await _run_agent_script(
                script, agent_id, user_id, timeout_seconds=90,
                name_prefix="__studio_session_open_", process_label="Studio: abrir sessão",
            )
            info = None
            for line in raw.splitlines():
                if line.startswith("__SESSION__:"):
                    try:
                        info = json.loads(line[len("__SESSION__:"):])
                    except Exception:
                        pass
                elif line.startswith("__SESSION_ERROR__:"):
                    raise Exception(line[len("__SESSION_ERROR__:"):].strip())
            if not info:
                raise Exception("Não foi possível abrir a sessão do navegador no agente.")
            sessions[session_name] = {**info, "agent_id": agent_id, "engine": engine}
            note = " — profile persistente reaproveitado, extensões instaladas continuam disponíveis" if profile_name else ""
            result = f"Sessão '{session_name}' aberta ({engine}){note}."
            _store(result, var, ctx)
            return result

        # demais ações exigem que a sessão já esteja aberta
        session = sessions.get(session_name)
        if not session:
            raise Exception(
                f"Sessão '{session_name}' não encontrada. "
                f"Adicione um passo 'Abrir sessão' com esse nome antes deste."
            )

        if t == "browser_close":
            script = _gen_session_close_script(session["pid"], session["port"], session["user_data_dir"],
                                                engine=session.get("engine", "playwright"),
                                                persistent=session.get("persistent", False))
            await _run_agent_script(
                script, session.get("agent_id", agent_id), user_id, timeout_seconds=60,
                name_prefix="__studio_session_close_", process_label="Studio: fechar sessão",
            )
            sessions.pop(session_name, None)
            result = f"Sessão '{session_name}' encerrada."
            _store(result, var, ctx)
            return result

        action_type = t[len("browser_"):]  # click | type | extract | wait | screenshot | captcha_*
        script = _gen_session_action_script(action_type, session["port"], cfg.get("target", ""), cfg.get("value", ""),
                                             ctx, engine=session.get("engine", "playwright"))
        # captcha_wait espera alguém resolver na tela — pode legitimamente demorar mais
        # que as demais ações (padrão 120s configurável no campo 'value'); o timeout de
        # polling aqui precisa ser maior que o deadline interno do script, senão a HAC
        # desiste antes do próprio script.
        if action_type == "captcha_wait":
            dispatch_timeout = int(float(cfg.get("value") or 120)) + 30
        else:
            dispatch_timeout = 60
        raw = await _run_agent_script(
            script, session.get("agent_id", agent_id), user_id, timeout_seconds=dispatch_timeout,
            name_prefix="__studio_session_act_", process_label="Studio: ação na sessão",
        )
        result = raw.strip()
        _store(result, var, ctx)
        return result

    return f"[{t}] step executado"


# ── Execução no agente (worker) em vez do servidor ─────────────────────────
#
# Todo tipo de step PODE ser configurado (campo run_on='agent') para rodar na
# máquina de um agente worker em vez de dentro do processo da API — útil para
# ações de arquivo/mídia que precisam enxergar o filesystem local do usuário,
# não o do servidor na nuvem. Os únicos tipos que NUNCA podem rodar no agente
# são os que dependem de acesso direto ao Mongo (call_ai_agent, call_pipeline,
# call_automation) ou que já têm seu próprio mecanismo de despacho pro agente
# (browser*) — para esses, o campo run_on é ignorado.
_AGENT_EXCLUDED_TYPES = {
    "call_ai_agent", "call_pipeline", "call_automation",
    "browser", "browser_open", "browser_click", "browser_type",
    "browser_extract", "browser_wait", "browser_screenshot", "browser_close",
    "browser_captcha_detect", "browser_captcha_wait", "browser_captcha_solve_image",
}


def _gen_local_step_script(step: dict, ctx: dict) -> str:
    """Gera um script Python autocontido que executa UM step localmente na
    máquina do agente, reutilizando o mesmo código de `_exec_step` (via
    inspect.getsource — sem duplicar a lógica das ~136 ações elegíveis).
    Os ramos que dependem do Mongo (call_ai_agent/call_pipeline/call_automation)
    ou de sessão de navegador (browser*) ficam presentes no texto mas nunca são
    alcançados, pois esses tipos nunca chegam aqui (ver _AGENT_EXCLUDED_TYPES)."""
    payload = {
        "step": {"type": step["type"], "config": step.get("config", {})},
        "ctx": {
            "input": ctx.get("input", ""),
            "output": ctx.get("output", ""),
            "vars": ctx.get("vars", {}),
        },
    }
    blob = base64.b64encode(json.dumps(payload, ensure_ascii=False, default=str).encode("utf-8")).decode("ascii")

    # Só stdlib garantido vai no topo sem proteção. 'httpx' é pacote pip de
    # terceiro (usado por muitas ações via `import httpx` no topo do módulo
    # original) — importar ele incondicionalmente aqui faria o script quebrar
    # ANTES de chegar no bloco que tenta instalar sozinho, caso o agente não
    # o tenha. Por isso o import dele fica dentro de _run_once(), coberto
    # pelo mesmo try/except-e-reinstala que protege os outros pacotes.
    header = [
        "import sys, os, json, re, asyncio, base64, hashlib, io, secrets, shutil, subprocess, time, zipfile, unicodedata",
        "sys.stdout.reconfigure(encoding='utf-8', errors='replace')",
        "import glob as glob_module",
        "from datetime import datetime, timedelta",
        "",
        f'_payload = json.loads(base64.b64decode("{blob}").decode("utf-8"))',
        "",
    ]

    # _is_cancelled é chamado por dentro dos steps 'wait'/'random_wait' — se não for
    # embutido aqui, um "Aguardar" configurado pra rodar no agente quebra com NameError.
    # No agente, ctx nunca tem _run_id nem _disconnect_check, então a função só retorna
    # False de cara (o cancelamento de um step rodando no agente é tratado por fora,
    # matando o processo do lado do worker — aqui é só pra não faltar o nome).
    body_parts = [
        inspect.getsource(_resolve_var_expr),
        inspect.getsource(_sub),
        inspect.getsource(_store),
        inspect.getsource(_slugify_var),
        inspect.getsource(_find_tesseract_binary),
        inspect.getsource(_tesseract_installed),
        inspect.getsource(_find_poppler_bin_dir),
        inspect.getsource(_poppler_installed),
        inspect.getsource(_ensure_native_binary),
        inspect.getsource(_pix_tlv),
        inspect.getsource(_pix_crc16),
        inspect.getsource(_build_pix_payload),
        f"_MODULE_TO_PACKAGE = {_MODULE_TO_PACKAGE!r}\n",
        f"_NO_AUTO_INSTALL = {_NO_AUTO_INSTALL!r}\n",
        inspect.getsource(_pip_package_for_module),
        inspect.getsource(_try_pip_install),
        inspect.getsource(_is_cancelled),
        inspect.getsource(_exec_step),
    ]

    footer = [
        "async def _run_once(step, ctx):",
        "    global httpx",
        "    import httpx as _httpx_mod",
        "    httpx = _httpx_mod",
        "    return await _exec_step(step, ctx)",
        "",
        "async def _main():",
        "    step = _payload['step']",
        "    ctx = _payload['ctx']",
        "    ctx.setdefault('vars', {})",
        "    try:",
        "        try:",
        "            output = await _run_once(step, ctx)",
        "        except ModuleNotFoundError as _e:",
        "            _missing = _e.name or ''",
        "            _pkg = _pip_package_for_module(_missing)",
        "            _ok, _msg = await _try_pip_install(_pkg)",
        "            if not _ok:",
        "                raise Exception(\"Pacote '\" + _pkg + \"' ausente no agente (interpretador \" + sys.executable + \") e nao pode ser instalado: \" + _msg)",
        "            try:",
        "                output = await _run_once(step, ctx)",
        "            except ModuleNotFoundError as _e2:",
        "                raise Exception(\"Instalei o pacote '\" + _pkg + \"' no agente (pip reportou sucesso), mas o modulo '\" + _missing + \"' continua ausente. Costuma acontecer quando ha mais de um Python nessa maquina (o pip instalou num, mas quem esta rodando e outro) ou quando falta algo alem do pip install. Interpretador usado: \" + sys.executable + \". Tente rodar '\" + sys.executable + \" -m pip install \" + _pkg + \"' manualmente no agente pra ver o erro real.\")",
        "        print('__STEP_RESULT__:' + json.dumps({'ok': True, 'output': str(output), 'vars': ctx['vars']}, ensure_ascii=False, default=str))",
        "    except Exception as _e:",
        "        print('__STEP_RESULT__:' + json.dumps({'ok': False, 'error': str(_e)}, ensure_ascii=False))",
        "",
        "asyncio.run(_main())",
    ]

    return "\n".join(header) + "\n\n" + "\n\n".join(body_parts) + "\n\n" + "\n".join(footer)


async def _exec_step_via_agent(step: dict, ctx: dict) -> str:
    """Despacha um step para rodar na máquina de um agente worker, aguarda o
    resultado e mescla as variáveis atualizadas de volta no contexto."""
    agent_id = ctx.get("agent_id", "")
    user_id = ctx.get("user_id", "")
    if not agent_id:
        raise Exception(
            "Este passo está configurado para rodar no agente ('Onde executar: Agente'), "
            "mas nenhum agente foi selecionado. Escolha um agente no seletor ao lado do "
            "botão Executar antes de salvar."
        )
    script = _gen_local_step_script(step, ctx)
    raw = await _run_agent_script(
        script, agent_id, user_id, timeout_seconds=90,
        name_prefix="__studio_local_step_", process_label=f"Studio: {step.get('name') or step['type']} (no agente)",
    )
    parsed = None
    for line in raw.splitlines():
        if line.startswith("__STEP_RESULT__:"):
            try:
                parsed = json.loads(line[len("__STEP_RESULT__:"):])
            except Exception:
                pass
    if parsed is None:
        raise Exception(f"O agente não retornou um resultado reconhecível. Saída bruta: {raw[:500]}")
    if not parsed.get("ok"):
        raise Exception(parsed.get("error") or "Falhou no agente sem mensagem de erro")
    ctx["vars"].update(parsed.get("vars") or {})
    return parsed.get("output", "")


_background_tasks: set = set()


def _spawn_background(coro):
    """Dispara uma corrotina em background sem esperar o resultado, mantendo uma
    referência forte até ela terminar (senão o garbage collector pode derrubar a
    task no meio, um problema clássico do asyncio.create_task solto)."""
    task = asyncio.create_task(coro)
    _background_tasks.add(task)
    task.add_done_callback(_background_tasks.discard)
    return task


async def _push_live(ctx: dict, result: dict):
    """Registra um resultado de step tanto na lista local (via chamador) quanto,
    se a execução tiver um run_id associado, grava no Mongo na hora — é isso que
    permite o frontend acompanhar o progresso em tempo real via polling, em vez
    de só ver tudo de uma vez quando a automação inteira termina."""
    live = ctx.setdefault("_live_results", [])
    live.append(result)
    run_id = ctx.get("_run_id")
    if run_id:
        try:
            await studio_runs_col.update_one(
                {"_id": run_id},
                {"$set": {"steps_result": live, "output": ctx.get("output", "")}},
            )
        except Exception:
            pass  # atualização de progresso nunca deve derrubar a execução


async def _is_cancelled(ctx: dict) -> bool:
    disconnect_check = ctx.get("_disconnect_check")
    if disconnect_check and await disconnect_check():
        return True
    run_id = ctx.get("_run_id")
    if run_id:
        doc = await studio_runs_col.find_one({"_id": run_id}, {"cancel_requested": 1})
        if doc and doc.get("cancel_requested"):
            return True
    return False


async def _exec_step_list(steps: list, ctx: dict) -> tuple:
    """Executa recursivamente uma lista de steps. Retorna (results, failed)."""
    results = []
    i = 0
    while i < len(steps):
        if await _is_cancelled(ctx):
            ctx["_cancelled"] = True
            cancelled_result = {
                "step_id": "", "step_name": "Execução cancelada", "step_type": "",
                "status": "cancelled", "output": "", "error": "Cancelada pelo usuário",
                "duration_ms": 0, "condition_result": None,
            }
            results.append(cancelled_result)
            await _push_live(ctx, cancelled_result)
            return results, True

        step = steps[i]
        step_type = step.get("type", "") if isinstance(step, dict) else step.type
        cfg = (step.get("config", {}) if isinstance(step, dict) else step.config.__dict__) or {}
        step_id = step.get("id", "") if isinstance(step, dict) else step.id
        step_name = (step.get("name") if isinstance(step, dict) else step.name) or f"Passo {i + 1}"

        enabled = step.get("enabled", True) if isinstance(step, dict) else getattr(step, "enabled", True)
        if not enabled:
            skip_result = {
                "step_id": step_id, "step_name": step_name, "step_type": step_type,
                "status": "skipped", "output": "[desabilitado]", "error": "",
                "duration_ms": 0, "condition_result": None,
            }
            results.append(skip_result)
            await _push_live(ctx, skip_result)
            i += 1
            continue

        step_start = time.time()

        result = {
            "step_id": step_id,
            "step_name": step_name,
            "step_type": step_type,
            "status": "success",
            "output": "",
            "error": "",
            "duration_ms": 0,
            "condition_result": None,
        }

        try:
            if step_type == "condition":
                cond_result = _eval_condition(
                    ctx["output"],
                    cfg.get("operator", "contains"),
                    cfg.get("condition_value", ""),
                )
                result["condition_result"] = cond_result
                result["output"] = f"Condição: {'VERDADEIRO ✓' if cond_result else 'FALSO ✗'} ({cfg.get('operator')} '{cfg.get('condition_value')}')"
                result["duration_ms"] = int((time.time() - step_start) * 1000)
                results.append(result)
                await _push_live(ctx, result)

                has_children = (
                    "children_true" in (step if isinstance(step, dict) else {})
                    or "children_false" in (step if isinstance(step, dict) else {})
                    or (not isinstance(step, dict) and (step.children_true or step.children_false))
                )
                if has_children:
                    # New children-based branching
                    if isinstance(step, dict):
                        branch = step.get("children_true", []) if cond_result else step.get("children_false", [])
                    else:
                        branch = step.children_true if cond_result else step.children_false
                    branch_results, branch_failed = await _exec_step_list(branch, ctx)
                    results.extend(branch_results)
                    if branch_failed:
                        return results, True
                    if ctx.get("_break_loop"):
                        return results, False
                    i += 1
                else:
                    # Legacy else_step_id jump logic
                    if not cond_result:
                        else_id = cfg.get("else_step_id", "")
                        if not else_id:
                            return results, False
                        else:
                            i = next((idx for idx, s in enumerate(steps)
                                      if (s.get("id") if isinstance(s, dict) else s.id) == else_id),
                                     len(steps))
                    else:
                        i += 1
                continue

            elif step_type == "loop_count":
                count = int(cfg.get("count", 3))
                idx_var = cfg.get("index_variable", "loop_index")
                children = (step.get("children", []) if isinstance(step, dict) else step.children) or []
                result["output"] = f"Loop: {count} iteração(ões), variável '{idx_var}'"
                result["duration_ms"] = int((time.time() - step_start) * 1000)
                results.append(result)
                await _push_live(ctx, result)

                ctx["_loop_depth"] = ctx.get("_loop_depth", 0) + 1
                try:
                    for iter_i in range(count):
                        ctx["vars"][idx_var] = str(iter_i)
                        iter_results, iter_failed = await _exec_step_list(children, ctx)
                        for r in iter_results:
                            r = {**r, "step_name": f"[{iter_i + 1}/{count}] {r['step_name']}"}
                            results.append(r)
                        if iter_failed:
                            return results, True
                        if ctx.get("_break_loop"):
                            ctx["_break_loop"] = False
                            break
                finally:
                    ctx["_loop_depth"] -= 1
                i += 1
                continue

            elif step_type == "foreach":
                raw_list = _sub(cfg.get("list_source", "{output}"), ctx)
                try:
                    items = json.loads(raw_list)
                    if not isinstance(items, list):
                        items = [items]
                except Exception:
                    items = [x for x in raw_list.split("\n") if x.strip()]
                item_var = cfg.get("item_variable", "item") or "item"
                children = (step.get("children", []) if isinstance(step, dict) else step.children) or []
                result["output"] = f"Foreach: {len(items)} item(ns), variável '{item_var}'"
                result["duration_ms"] = int((time.time() - step_start) * 1000)
                results.append(result)
                await _push_live(ctx, result)

                ctx["_loop_depth"] = ctx.get("_loop_depth", 0) + 1
                try:
                    for iter_i, item in enumerate(items):
                        ctx["vars"][item_var] = item if isinstance(item, str) else json.dumps(item, ensure_ascii=False)
                        ctx["vars"][f"{item_var}_index"] = str(iter_i)
                        iter_results, iter_failed = await _exec_step_list(children, ctx)
                        for r in iter_results:
                            r = {**r, "step_name": f"[{iter_i + 1}/{len(items)}] {r['step_name']}"}
                            results.append(r)
                        if iter_failed:
                            return results, True
                        if ctx.get("_break_loop"):
                            ctx["_break_loop"] = False
                            break
                finally:
                    ctx["_loop_depth"] -= 1
                i += 1
                continue

            elif step_type == "while_condition":
                children = (step.get("children", []) if isinstance(step, dict) else step.children) or []
                operator = cfg.get("operator", "not_empty")
                cond_value = cfg.get("condition_value", "")
                max_iter = int(cfg.get("max_iterations", 100) or 100)
                result["output"] = f"While: {operator} '{cond_value}' (máx {max_iter} iterações)"
                result["duration_ms"] = int((time.time() - step_start) * 1000)
                results.append(result)
                await _push_live(ctx, result)

                ctx["_loop_depth"] = ctx.get("_loop_depth", 0) + 1
                try:
                    iter_i = 0
                    while _eval_condition(ctx["output"], operator, cond_value) and iter_i < max_iter:
                        iter_results, iter_failed = await _exec_step_list(children, ctx)
                        for r in iter_results:
                            r = {**r, "step_name": f"[{iter_i + 1}] {r['step_name']}"}
                            results.append(r)
                        if iter_failed:
                            return results, True
                        if ctx.get("_break_loop"):
                            ctx["_break_loop"] = False
                            break
                        iter_i += 1
                finally:
                    ctx["_loop_depth"] -= 1
                i += 1
                continue

            elif step_type == "try_catch":
                try_children = (step.get("children_true", []) if isinstance(step, dict) else step.children_true) or []
                catch_children = (step.get("children_false", []) if isinstance(step, dict) else step.children_false) or []
                result["output"] = "Try/Catch"
                result["duration_ms"] = int((time.time() - step_start) * 1000)
                results.append(result)
                await _push_live(ctx, result)

                try_results, try_failed = await _exec_step_list(try_children, ctx)
                results.extend(try_results)
                if try_failed and not ctx.get("_break_loop"):
                    last_err = next((r["error"] for r in reversed(try_results) if r.get("error")), "erro desconhecido")
                    ctx["vars"]["error"] = last_err
                    catch_results, catch_failed = await _exec_step_list(catch_children, ctx)
                    results.extend(catch_results)
                    if catch_failed:
                        return results, True
                if ctx.get("_break_loop"):
                    return results, False
                i += 1
                continue

            elif step_type == "break_loop":
                inside_loop = ctx.get("_loop_depth", 0) > 0
                result["output"] = "Loop interrompido" if inside_loop else "Nenhum loop em execução — passo ignorado"
                result["status"] = "success" if inside_loop else "skipped"
                result["duration_ms"] = int((time.time() - step_start) * 1000)
                results.append(result)
                await _push_live(ctx, result)
                if inside_loop:
                    ctx["_break_loop"] = True
                    return results, False
                i += 1
                continue

            elif step_type == "parallel":
                children = (step.get("children", []) if isinstance(step, dict) else step.children) or []
                result["output"] = f"Paralelo: {len(children)} ação(ões) concorrentes"
                result["duration_ms"] = int((time.time() - step_start) * 1000)
                results.append(result)
                await _push_live(ctx, result)

                async def _run_one_parallel(child, base_ctx):
                    child_ctx = {**base_ctx, "vars": dict(base_ctx["vars"])}
                    r, f = await _exec_step_list([child], child_ctx)
                    return r, f, child_ctx

                outcomes = await asyncio.gather(*[_run_one_parallel(c, ctx) for c in children])
                any_failed = False
                outputs = []
                for r, f, child_ctx in outcomes:
                    results.extend(r)
                    if f:
                        any_failed = True
                    outputs.append(child_ctx.get("output", ""))
                    ctx["vars"].update(child_ctx["vars"])
                ctx["output"] = "\n".join(outputs)
                if any_failed:
                    return results, True
                i += 1
                continue

            else:
                step_dict = step if isinstance(step, dict) else step.model_dump()
                cfg_dict = step_dict.get("config", {}) or {}
                run_on_agent = cfg_dict.get("run_on") == "agent" and step_type not in _AGENT_EXCLUDED_TYPES
                if run_on_agent:
                    output = await _exec_step_via_agent(step_dict, ctx)
                else:
                    try:
                        output = await _exec_step(step_dict, ctx)
                    except ModuleNotFoundError as e:
                        missing_module = e.name or ""
                        if missing_module.split(".")[0] in _NO_AUTO_INSTALL:
                            raise
                        package = _pip_package_for_module(missing_module)
                        install_ok, install_msg = await _try_pip_install(package)
                        if not install_ok:
                            raise Exception(
                                f"Esta ação precisa do pacote Python '{package}' (módulo '{missing_module}'), "
                                f"que não está instalado e não pôde ser instalado automaticamente: {install_msg}. "
                                f"Adicione '{package}' ao requirements.txt e faça redeploy. "
                                f"(interpretador: {sys.executable})"
                            ) from e
                        # reinstalado com sucesso — tenta a ação de novo, uma única vez
                        try:
                            output = await _exec_step(step_dict, ctx)
                        except ModuleNotFoundError as e2:
                            raise Exception(
                                f"Instalei o pacote '{package}' (pip reportou sucesso), mas o módulo "
                                f"'{missing_module}' continua não encontrado. Isso costuma acontecer quando há "
                                f"mais de um Python instalado na máquina (o pip instalou num, mas quem está "
                                f"rodando é outro) ou quando o pacote precisa de mais alguma coisa além do pip "
                                f"install. Interpretador usado: {sys.executable}. Tente rodar "
                                f"'{sys.executable} -m pip install {package}' manualmente para ver o erro real."
                            ) from e2
                        output = f"[pacote '{package}' instalado automaticamente] {output}"
                result["output"] = str(output)[:3000]

        except Exception as e:
            result["status"] = "failed"
            result["error"] = str(e)
            result["duration_ms"] = int((time.time() - step_start) * 1000)
            results.append(result)
            await _push_live(ctx, result)
            return results, True

        result["duration_ms"] = int((time.time() - step_start) * 1000)
        results.append(result)
        await _push_live(ctx, result)
        i += 1

    return results, False


async def _close_remaining_sessions(ctx: dict):
    """Rede de segurança: ao final de uma execução (sucesso, falha ou exceção),
    fecha qualquer sessão de navegador que ainda esteja registrada em ctx['sessions']
    — evita deixar processos de navegador órfãos na máquina do agente quando o
    fluxo termina sem passar por um passo 'Fechar sessão' (ex: erro no meio)."""
    sessions = ctx.get("sessions") or {}
    if not sessions:
        return
    user_id = ctx.get("user_id", "")
    for name, session in list(sessions.items()):
        try:
            script = _gen_session_close_script(session["pid"], session["port"], session["user_data_dir"],
                                                engine=session.get("engine", "playwright"),
                                                persistent=session.get("persistent", False))
            await _run_agent_script(
                script, session.get("agent_id", ctx.get("agent_id", "")), user_id, timeout_seconds=60,
                name_prefix="__studio_session_close_", process_label="Studio: fechar sessão (limpeza)",
            )
        except Exception:
            pass
        sessions.pop(name, None)


async def _execute_automation(automation: dict, initial_input: str, trigger_type: str = "manual",
                               request: Request = None, run_id: str = None) -> dict:
    """Executa a automação do início ao fim. Se `run_id` for informado, assume que o
    chamador já inseriu o run_doc (status 'running') no Mongo — usado pelo endpoint
    manual, que cria o doc e retorna na hora, disparando a execução real em background,
    para o frontend poder acompanhar o progresso via polling. Sem `run_id`, o próprio
    doc é criado aqui (comportamento síncrono de sempre — usado por webhook e cron)."""
    started_at = datetime.utcnow()
    auto_created = run_id is None
    if auto_created:
        run_id = str(ObjectId())
        run_doc = {
            "_id": run_id,
            "automation_id": automation["_id"],
            "automation_name": automation["name"],
            "user_id": automation["user_id"],
            "trigger_type": trigger_type,
            "input": initial_input,
            "steps_result": [],
            "output": "",
            "status": "running",
            "started_at": started_at,
            "finished_at": None,
            "duration_ms": 0,
            "cancel_requested": False,
        }
        await studio_runs_col.insert_one(run_doc)
    else:
        run_doc = await studio_runs_col.find_one({"_id": run_id}) or {}
        started_at = run_doc.get("started_at", started_at)

    steps = automation.get("steps", [])
    ctx = {
        "input": initial_input, "output": initial_input, "vars": {}, "sessions": {},
        "agent_id": automation.get("agent_id", ""),
        "user_id": str(automation.get("user_id", "")),
        "_run_id": run_id,
    }
    if request is not None:
        ctx["_disconnect_check"] = request.is_disconnected

    try:
        steps_results, failed = await _exec_step_list(steps, ctx)
    finally:
        await _close_remaining_sessions(ctx)
    final_status = "cancelled" if ctx.get("_cancelled") else ("failed" if failed else "success")

    finished_at = datetime.utcnow()
    duration_ms = int((finished_at - started_at).total_seconds() * 1000)

    await studio_runs_col.update_one({"_id": run_id}, {"$set": {
        "steps_result": steps_results,
        "output": ctx["output"],
        "status": final_status,
        "finished_at": finished_at,
        "duration_ms": duration_ms,
    }})

    return {**run_doc, "steps_result": steps_results, "output": ctx["output"],
            "status": final_status, "finished_at": finished_at, "duration_ms": duration_ms}


# ─── CRUD ─────────────────────────────────────────────────────────

@router.post("", response_model=AutomationOut, status_code=201)
async def create_automation(body: AutomationCreate, request: Request, user: dict = Depends(get_current_user)):
    trigger = body.trigger.model_dump()
    if trigger.get("type") == "webhook" and not trigger.get("webhook_token"):
        trigger["webhook_token"] = secrets.token_urlsafe(20)
    now = datetime.utcnow()
    doc = {
        "_id": str(ObjectId()),
        "user_id": user["_id"],
        "name": body.name,
        "description": body.description,
        "trigger": trigger,
        "steps": [s.model_dump() for s in body.steps],
        "active": body.active,
        "agent_id": body.agent_id,
        "created_at": now,
        "updated_at": now,
    }
    await studio_automations_col.insert_one(doc)
    schedule = trigger.get("schedule", "") if trigger.get("type") == "cron" else ""
    await _upsert_linked_process(doc["_id"], user["_id"], body.name, body.description, body.agent_id, schedule)
    return _doc_to_out(doc, str(request.base_url).rstrip("/"))


@router.get("", response_model=List[AutomationOut])
async def list_automations(request: Request, user: dict = Depends(get_current_user)):
    base_url = str(request.base_url).rstrip("/")
    cursor = studio_automations_col.find({"user_id": user["_id"]}).sort("created_at", -1)
    return [_doc_to_out(doc, base_url) async for doc in cursor]


@router.get("/{automation_id}", response_model=AutomationOut)
async def get_automation(automation_id: str, request: Request, user: dict = Depends(get_current_user)):
    doc = await studio_automations_col.find_one({"_id": automation_id, "user_id": user["_id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Automação não encontrada")
    return _doc_to_out(doc, str(request.base_url).rstrip("/"))


@router.patch("/{automation_id}", response_model=AutomationOut)
async def update_automation(automation_id: str, body: AutomationUpdate, request: Request, user: dict = Depends(get_current_user)):
    doc = await studio_automations_col.find_one({"_id": automation_id, "user_id": user["_id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Automação não encontrada")
    data = body.model_dump(exclude_unset=True)
    if "trigger" in data and data["trigger"]:
        t = data["trigger"]
        if t.get("type") == "webhook" and not t.get("webhook_token"):
            t["webhook_token"] = doc.get("trigger", {}).get("webhook_token") or secrets.token_urlsafe(20)
    if "steps" in data and data["steps"] is not None:
        data["steps"] = [s if isinstance(s, dict) else s.model_dump() for s in data["steps"]]
    data["updated_at"] = datetime.utcnow()
    await studio_automations_col.update_one({"_id": automation_id}, {"$set": data})
    merged = {**doc, **data}
    trigger_merged = merged.get("trigger", {})
    schedule = trigger_merged.get("schedule", "") if trigger_merged.get("type") == "cron" else ""
    await _upsert_linked_process(
        automation_id, user["_id"],
        merged.get("name", doc["name"]),
        merged.get("description", doc.get("description", "")),
        merged.get("agent_id", doc.get("agent_id", "")),
        schedule,
    )
    return _doc_to_out(merged, str(request.base_url).rstrip("/"))


@router.delete("/{automation_id}", status_code=204)
async def delete_automation(automation_id: str, user: dict = Depends(get_current_user)):
    result = await studio_automations_col.delete_one({"_id": automation_id, "user_id": user["_id"]})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Automação não encontrada")
    await processes_col.delete_one({"studio_automation_id": automation_id, "user_id": user["_id"]})


class RunRequest(BaseModel):
    input: str = ""
    only_step_id: str = ""
    from_step_id: str = ""


@router.post("/{automation_id}/run", response_model=AutomationRunOut, status_code=202)
async def run_automation(automation_id: str, body: RunRequest, user: dict = Depends(get_current_user)):
    """Cria o registro da execução e retorna IMEDIATAMENTE (status 'running'), disparando
    a execução de verdade em background — o chamador acompanha o progresso via
    GET /{automation_id}/runs/{run_id}, feito a cada ~1s pelo frontend (polling), e o
    log preenche em tempo real conforme cada step termina, não tudo de uma vez no final.

    `only_step_id`/`from_step_id` (usados pelo menu de contexto do Builder — "Executar
    este passo" / "Executar a partir deste passo") só enxergam o NÍVEL PRINCIPAL do fluxo
    (não um passo dentro de um branch de condição/loop/try-catch), porque rodar um trecho
    aninhado fora do laço/condição que o envolve não teria contexto (variável de loop,
    resultado da condição etc.) pra fazer sentido sozinho."""
    doc = await studio_automations_col.find_one({"_id": automation_id, "user_id": user["_id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Automação não encontrada")

    steps = doc.get("steps", [])
    if body.only_step_id or body.from_step_id:
        target_id = body.only_step_id or body.from_step_id
        idx = next((i for i, s in enumerate(steps) if s.get("id") == target_id), None)
        if idx is None:
            raise HTTPException(status_code=400, detail="Passo não encontrado no nível principal do fluxo")
        if body.only_step_id:
            sliced = [dict(steps[idx], enabled=True)]
        else:
            sliced = steps[idx:]
        doc = {**doc, "steps": sliced}

    run_id = str(ObjectId())
    now = datetime.utcnow()
    run_doc = {
        "_id": run_id,
        "automation_id": automation_id,
        "automation_name": doc["name"],
        "user_id": doc["user_id"],
        "trigger_type": "manual",
        "input": body.input,
        "steps_result": [],
        "output": "",
        "status": "running",
        "started_at": now,
        "finished_at": None,
        "duration_ms": 0,
        "cancel_requested": False,
    }
    await studio_runs_col.insert_one(run_doc)
    _spawn_background(_execute_automation(doc, body.input, trigger_type="manual", run_id=run_id))
    return _run_doc_to_out(run_doc)


@router.get("/{automation_id}/runs/{run_id}", response_model=AutomationRunOut)
async def get_run(automation_id: str, run_id: str, user: dict = Depends(get_current_user)):
    doc = await studio_automations_col.find_one({"_id": automation_id, "user_id": user["_id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Automação não encontrada")
    run = await studio_runs_col.find_one({"_id": run_id, "automation_id": automation_id})
    if not run:
        raise HTTPException(status_code=404, detail="Execução não encontrada")
    return _run_doc_to_out(run)


@router.post("/{automation_id}/runs/{run_id}/cancel", status_code=204)
async def cancel_run(automation_id: str, run_id: str, user: dict = Depends(get_current_user)):
    doc = await studio_automations_col.find_one({"_id": automation_id, "user_id": user["_id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Automação não encontrada")
    result = await studio_runs_col.update_one(
        {"_id": run_id, "automation_id": automation_id, "status": "running"},
        {"$set": {"cancel_requested": True}},
    )
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Execução em andamento não encontrada")


@router.get("/{automation_id}/runs", response_model=List[AutomationRunOut])
async def list_runs(automation_id: str, user: dict = Depends(get_current_user)):
    doc = await studio_automations_col.find_one({"_id": automation_id, "user_id": user["_id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Automação não encontrada")
    cursor = studio_runs_col.find({"automation_id": automation_id}).sort("started_at", -1).limit(20)
    return [_run_doc_to_out(d) async for d in cursor]


# ─── TESTE (executa o que está no Builder sem precisar salvar) ────

_ADHOC_AUTOMATION_ID = "__adhoc__"


class AdHocRunRequest(BaseModel):
    steps: List[AutomationStep]
    input: str = ""
    agent_id: str = ""


@router.post("/test-run", response_model=AutomationRunOut, status_code=202)
async def test_run_adhoc(body: AdHocRunRequest, user: dict = Depends(get_current_user)):
    """Executa a lista de passos exatamente como está no Builder agora, sem exigir que a
    automação tenha sido salva antes — usado pelos botões Executar / Executar este passo /
    Executar a partir deste passo enquanto o usuário ainda está editando o fluxo."""
    doc = {
        "_id": _ADHOC_AUTOMATION_ID,
        "name": "(teste não salvo)",
        "steps": [s.model_dump() for s in body.steps],
        "user_id": user["_id"],
        "agent_id": body.agent_id,
    }
    run_id = str(ObjectId())
    now = datetime.utcnow()
    run_doc = {
        "_id": run_id,
        "automation_id": _ADHOC_AUTOMATION_ID,
        "automation_name": doc["name"],
        "user_id": user["_id"],
        "trigger_type": "manual",
        "input": body.input,
        "steps_result": [],
        "output": "",
        "status": "running",
        "started_at": now,
        "finished_at": None,
        "duration_ms": 0,
        "cancel_requested": False,
    }
    await studio_runs_col.insert_one(run_doc)
    _spawn_background(_execute_automation(doc, body.input, trigger_type="manual", run_id=run_id))
    return _run_doc_to_out(run_doc)


@router.get("/test-run/{run_id}", response_model=AutomationRunOut)
async def get_test_run_adhoc(run_id: str, user: dict = Depends(get_current_user)):
    run = await studio_runs_col.find_one(
        {"_id": run_id, "automation_id": _ADHOC_AUTOMATION_ID, "user_id": user["_id"]}
    )
    if not run:
        raise HTTPException(status_code=404, detail="Execução não encontrada")
    return _run_doc_to_out(run)


@router.post("/test-run/{run_id}/cancel", status_code=204)
async def cancel_test_run_adhoc(run_id: str, user: dict = Depends(get_current_user)):
    result = await studio_runs_col.update_one(
        {"_id": run_id, "automation_id": _ADHOC_AUTOMATION_ID, "user_id": user["_id"], "status": "running"},
        {"$set": {"cancel_requested": True}},
    )
    if result.matched_count == 0:
        raise HTTPException(status_code=404, detail="Execução em andamento não encontrada")


# ─── WEBHOOK ──────────────────────────────────────────────────────

@router.post("/webhook/{token}")
async def webhook_trigger(token: str, request: Request):
    doc = await studio_automations_col.find_one({"trigger.webhook_token": token, "active": True})
    if not doc:
        raise HTTPException(status_code=404, detail="Webhook não encontrado")
    try:
        body = await request.json()
        input_text = json.dumps(body)
    except Exception:
        raw = await request.body()
        input_text = raw.decode("utf-8", errors="replace")
    run = await _execute_automation(doc, input_text, trigger_type="webhook")
    return {"run_id": run["_id"], "status": run["status"], "output": run["output"]}
